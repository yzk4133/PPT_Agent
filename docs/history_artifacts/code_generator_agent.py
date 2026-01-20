"""代码生成Agent - 负责生成PPTX Python代码

将内容和设计组合，生成可执行的python-pptx代码
"""

from typing import Dict, Any
from .agent_base import BaseAgent, AgentConfig, AgentMessage


class CodeGeneratorAgent(BaseAgent):
    """代码生成Agent

    职责:
    - 将内容转换为PPTX代码
    - 应用设计规范
    - 处理模板和布局
    - 确保代码可执行性
    """

    def __init__(self, config: AgentConfig):
        super().__init__(config)
        self.template_cache = {}

    async def process(self, message: AgentMessage) -> AgentMessage:
        """处理代码生成请求"""
        task_type = message.content.get("task_type")

        if task_type == "generate_ppt_code":
            return await self._generate_ppt_code(message)
        elif task_type == "validate_code":
            return await self._validate_code(message)
        else:
            raise ValueError(f"未知的任务类型: {task_type}")

    async def _generate_ppt_code(self, message: AgentMessage) -> AgentMessage:
        """生成PPTX代码"""
        content_data = message.content.get("content_data", {})
        design_scheme = message.content.get("design_scheme", {})
        template = message.content.get("template", "default")

        # 生成导入和初始化代码
        init_code = self._generate_init_code()

        # 生成幻灯片代码
        slides_code = self._generate_slides_code(content_data, design_scheme)

        # 生成保存代码
        save_code = self._generate_save_code()

        # 组合完整代码
        full_code = f"{init_code}\n\n{slides_code}\n\n{save_code}"

        return self._create_response(
            to_agent=message.from_agent,
            content={
                "task_type": "code_generated",
                "code": full_code,
                "language": "python"
            }
        )

    def _generate_init_code(self) -> str:
        """生成初始化代码"""
        return '''from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)'''

    def _generate_slides_code(self, content: Dict, design: Dict) -> str:
        """生成幻灯片代码"""
        code_parts = []

        for i, slide in enumerate(content.get("slides", [])):
            slide_code = f'''
# 幻灯片 {i+1}: {slide.get('title', '未命名')}
slide_layout = prs.slide_layouts[{slide.get('layout', 1)}]
slide = prs.slides.add_slide(slide_layout)

# 设置标题
title = slide.shapes.title
title.text = "{slide.get('title', '')}"

# 设置内容
if slide.get('content'):
    content_box = slide.placeholders[{slide.get('content_placeholder', 1)}]
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for j, point in enumerate(slide.get('content', [])):
        if j == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = {slide.get('indent', 0)}
'''
            code_parts.append(slide_code)

        return "\n".join(code_parts)

    def _generate_save_code(self) -> str:
        """生成保存代码"""
        return '''# 保存演示文稿
output_path = "output/presentation.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")'''

    async def _validate_code(self, message: AgentMessage) -> AgentMessage:
        """验证代码的正确性"""
        code = message.content.get("code", "")

        # 执行语法检查
        try:
            compile(code, '<string>', 'exec')
            is_valid = True
            error = None
        except SyntaxError as e:
            is_valid = False
            error = str(e)

        return self._create_response(
            to_agent=message.from_agent,
            content={
                "task_type": "code_validated",
                "is_valid": is_valid,
                "error": error
            }
        )

    def _create_response(self, to_agent: str, content: Dict[str, Any]) -> AgentMessage:
        """创建响应消息"""
        return AgentMessage.create(
            from_agent=self.config.name,
            to_agent=to_agent,
            content=content
        )
