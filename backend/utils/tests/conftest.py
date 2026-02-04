"""
Shared fixtures for utils tests
"""
import pytest
from unittest.mock import MagicMock
from pptx.util import Pt, Inches
from io import BytesIO
import sys
import os

# Add parent directory to path for imports

@pytest.fixture
def sample_slide_xml():
    """示例幻灯片XML"""
    return {
        "simple": '<SECTION layout="vertical"><H1>标题</H1><P>内容</P></SECTION>',
        "with_bullets": '''<SECTION>
            <H1>要点列表</H1>
            <BULLETS>
                <P>要点一</P>
                <P>要点二</P>
                <P>要点三</P>
            </BULLETS>
        </SECTION>''',
        "with_image": '''<SECTION>
            <H1>图片页</H1>
            <IMG src="http://example.com/img.jpg"/>
        </SECTION>''',
        "with_columns": '''<SECTION layout="vertical">
            <H1>分栏内容</H1>
            <COLUMNS>
                <P>第一列内容</P>
                <P>第二列内容</P>
            </COLUMNS>
        </SECTION>''',
        "empty": '<SECTION layout="vertical"></SECTION>',
        "no_title": '<SECTION layout="vertical"><P>只有内容没有标题</P></SECTION>',
    }

@pytest.fixture
def sample_ppt_json():
    """示例PPT JSON数据"""
    return {
        "title": "测试演示文稿",
        "sections": [
            {
                "id": "section1",
                "content": [
                    {"type": "h1", "children": [{"text": "第一章"}]},
                    {"type": "bullets", "children": [
                        {"type": "bullet", "children": [
                            {"type": "h3", "children": [{"text": "要点1"}]},
                            {"type": "p", "children": [{"text": "详细内容1"}]}
                        ]},
                        {"type": "bullet", "children": [
                            {"type": "h3", "children": [{"text": "要点2"}]},
                            {"type": "p", "children": [{"text": "详细内容2"}]}
                        ]}
                    ]}
                ],
                "rootImage": {"url": "http://example.com/img.jpg", "alt": "测试图", "background": False}
            },
            {
                "id": "section2",
                "content": [
                    {"type": "h1", "children": [{"text": "第二章"}]},
                    {"type": "p", "children": [{"text": "段落内容"}]}
                ]
            }
        ],
        "references": ["参考文献1", "参考文献2", "参考文献3"]
    }

@pytest.fixture
def mock_presentation():
    """Mock PowerPoint Presentation"""
    pres = MagicMock()
    pres.slide_width = 9144000  # 10英寸
    pres.slide_height = 6858000  # 7.5英寸
    pres.slide_layouts = [MagicMock(name=f"layout{i}") for i in range(31)]

    # Mock slides collection
    pres.slides = MagicMock()
    pres.slides.__len__ = MagicMock(return_value=0)

    return pres

@pytest.fixture
def mock_slide():
    """Mock Slide对象"""
    slide = MagicMock()
    slide.shapes = []

    # Mock文本框形状
    def create_text_shape(shape_id, name):
        text_shape = MagicMock()
        text_shape.shape_id = shape_id
        text_shape.has_text_frame = True
        text_shape.name = name
        text_shape.is_placeholder = False
        text_shape.text = ""

        text_frame = MagicMock()
        text_frame.word_wrap = True
        text_frame.auto_size = 1  # TEXT_TO_FIT_SHAPE
        text_frame.vertical_anchor = 3  # MIDDLE
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.paragraphs = [MagicMock()]

        p = text_frame.paragraphs[0]
        p.alignment = 0
        p.font = MagicMock()
        p.font.size = Pt(18)

        text_shape.text_frame = text_frame
        return text_shape

    # 添加标题形状
    slide.shapes.append(create_text_shape(2, "Title"))
    # 添加内容形状
    slide.shapes.append(create_text_shape(3, "Content"))

    slide.slide_layout = MagicMock()
    slide.slide_layout.name = "Test Layout"

    return slide

@pytest.fixture
def mock_shape():
    """Mock Shape对象"""
    shape = MagicMock()
    shape.shape_id = 1
    shape.has_text_frame = True
    shape.name = "Test Shape"
    shape.is_placeholder = False
    shape.text = ""

    text_frame = MagicMock()
    text_frame.word_wrap = True
    text_frame.auto_size = 1
    text_frame.vertical_anchor = 3
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.paragraphs = [MagicMock()]

    p = text_frame.paragraphs[0]
    p.alignment = 0
    p.font = MagicMock()
    p.font.size = Pt(18)

    shape.text_frame = text_frame
    shape.width = Inches(5)
    shape.height = Inches(3)
    shape.left = Inches(1)
    shape.top = Inches(1)

    return shape
