"""
Tests for slide strategies in ppt_generator module
"""
import pytest
from unittest.mock import MagicMock, Mock, patch
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
import sys
import os

# Add parent directory to path

from utils.save_ppt.strategies import (
    SlideStrategy,
    TitleSlideStrategy,
    ContentSlideStrategy,
    TableOfContentsSlideStrategy,
    ImageSlideStrategy,
    SubSectionSlideStrategy,
    ReferencesSlideStrategy,
    EndSlideStrategy,
)
from utils.save_ppt.config import SlideConfig

@pytest.fixture
def mock_presentation():
    """Mock Presentation对象"""
    pres = MagicMock()
    pres.slide_width = 9144000  # 10英寸
    pres.slide_height = 6858000  # 7.5英寸
    pres.slide_layouts = [MagicMock(name=f"layout{i}") for i in range(31)]
    pres.slides = MagicMock()
    return pres

@pytest.fixture
def mock_slide():
    """Mock Slide对象"""
    slide = MagicMock()

    # Mock文本框形状
    def create_text_shape(shape_id, name):
        text_shape = MagicMock()
        text_shape.shape_id = shape_id
        text_shape.has_text_frame = True
        text_shape.name = name
        text_shape.is_placeholder = False
        text_shape.text = ""

        text_frame = MagicMock()
        text_frame.word_wrap = True
        text_frame.auto_size = 1
        text_frame.vertical_anchor = 3
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.paragraphs = [MagicMock()]

        p = text_frame.paragraphs[0]
        p.alignment = 0
        p.font = MagicMock()
        p.font.size = Pt(18)

        text_shape.text_frame = text_frame
        return text_shape

    slide.shapes = [create_text_shape(2, "Title"), create_text_shape(9, "Date")]

    # 为每个shape设置位置和尺寸属性（避免格式化错误）
    for shape in slide.shapes:
        shape.left = 914400  # 1英寸
        shape.top = 914400   # 1英寸
        shape.width = 914400  # 1英寸
        shape.height = 914400  # 1英寸

    slide.slide_layout = MagicMock()
    slide.slide_layout.name = "Test Layout"
    slide.slide_layout.index = lambda: 0

    return slide

@pytest.fixture
def mock_image_slide():
    """Mock for image slide"""
    slide = MagicMock()

    # Mock图片占位符
    img_placeholder = MagicMock()
    img_placeholder.shape_id = 3
    img_placeholder.has_text_frame = False
    img_placeholder.name = "Image Placeholder"
    img_placeholder.width = Inches(6)
    img_placeholder.height = Inches(4)
    img_placeholder.left = Inches(2)
    img_placeholder.top = Inches(1.5)

    # 文本形状
    title_shape = MagicMock()
    title_shape.shape_id = 2
    title_shape.has_text_frame = True
    title_shape.name = "Title"
    title_shape.is_placeholder = False
    title_shape.text = ""

    desc_shape = MagicMock()
    desc_shape.shape_id = 4
    desc_shape.has_text_frame = True
    desc_shape.name = "Description"
    desc_shape.is_placeholder = False
    desc_shape.text = ""

    slide.shapes = [img_placeholder, title_shape, desc_shape]

    # 为每个shape设置位置和尺寸属性
    for shape in slide.shapes:
        shape.left = 914400  # 1英寸
        shape.top = 914400   # 1英寸
        shape.width = 914400  # 1英寸
        shape.height = 914400  # 1英寸

    return slide

@pytest.fixture
def config():
    """SlideConfig fixture"""
    return SlideConfig()

def create_mock_slide_with_shapes(shape_configs):
    """
    创建带有多个shape的mock slide

    Args:
        shape_configs: 列表，每项是 {shape_id, name, has_text_frame}

    Returns:
        Mock slide对象
    """
    slide = MagicMock()

    shapes = []
    for config in shape_configs:
        shape = MagicMock()
        shape.shape_id = config["shape_id"]
        shape.has_text_frame = config.get("has_text_frame", True)
        shape.name = config["name"]
        shape.is_placeholder = False
        shape.text = ""

        if shape.has_text_frame:
            text_frame = MagicMock()
            text_frame.word_wrap = True
            text_frame.auto_size = 1
            text_frame.vertical_anchor = 3
            text_frame.margin_left = 0.1
            text_frame.margin_right = 0.1
            text_frame.margin_top = 0.05
            text_frame.margin_bottom = 0.05
            text_frame.paragraphs = [MagicMock()]
            text_frame.paragraphs[0].alignment = 0
            text_frame.paragraphs[0].font = MagicMock()
            text_frame.paragraphs[0].font.size = 18
            shape.text_frame = text_frame

        # 设置位置和尺寸属性
        shape.left = 914400  # 1英寸
        shape.top = 914400   # 1英寸
        shape.width = 914400  # 1英寸
        shape.height = 914400  # 1英寸

        shapes.append(shape)

    slide.shapes = shapes
    slide.slide_layout = MagicMock()
    slide.slide_layout.name = "Test Layout"
    slide.slide_layout.index = lambda: 0

    return slide

def setup_mock_slide_position_attributes(slide):
    """
    为slide中的所有shapes设置位置和尺寸属性
    这是辅助函数，用于避免格式化错误
    """
    for shape in slide.shapes:
        if not hasattr(shape, 'left') or shape.left is None:
            shape.left = 914400  # 1英寸
        if not hasattr(shape, 'top') or shape.top is None:
            shape.top = 914400   # 1英寸
        if not hasattr(shape, 'width') or shape.width is None:
            shape.width = 914400  # 1英寸
        if not hasattr(shape, 'height') or shape.height is None:
            shape.height = 914400  # 1英寸
    return slide

class TestSlideStrategy:
    """测试SlideStrategy基类"""

    def test_slide_strategy_is_abstract(self):
        """测试SlideStrategy是抽象类，不能直接实例化"""
        from abc import ABC
        assert issubclass(SlideStrategy, ABC)
        # SlideStrategy有抽象方法create_slide，不能直接实例化

class TestTitleSlideStrategy:
    """测试TitleSlideStrategy"""

    @patch('utils.save_ppt.strategies.title_slide.datetime')
    def test_create_slide_basic(self, mock_datetime, mock_presentation, mock_slide, config):
        """测试基本标题页创建"""
        # Mock日期
        mock_datetime.datetime.now.return_value.strftime.return_value = "2025年01月01日"

        with patch.object(mock_presentation.slides, 'add_slide', return_value=mock_slide):
            strategy = TitleSlideStrategy(mock_presentation, config)
            strategy.create_slide("测试标题")

            assert strategy.slide_counter == 1

    @patch('utils.save_ppt.strategies.title_slide.datetime')
    def test_create_slide_with_long_title(self, mock_datetime, mock_presentation, mock_slide, config):
        """测试长标题"""
        mock_datetime.datetime.now.return_value.strftime.return_value = "2025年01月01日"

        with patch.object(mock_presentation.slides, 'add_slide', return_value=mock_slide):
            strategy = TitleSlideStrategy(mock_presentation, config)
            long_title = "这是一个非常非常长的标题" * 10
            strategy.create_slide(long_title)

            assert strategy.slide_counter == 1

class TestContentSlideStrategy:
    """测试ContentSlideStrategy"""

    def test_create_slide_basic(self, mock_presentation, mock_slide, config):
        """测试基本内容页创建"""
        with patch.object(mock_presentation.slides, 'add_slide', return_value=mock_slide):
            strategy = ContentSlideStrategy(mock_presentation, config)

            # 创建多个文本形状
            def create_text_shape(shape_id, name):
                shape = MagicMock()
                shape.shape_id = shape_id
                shape.has_text_frame = True
                shape.name = name
                shape.is_placeholder = False
                shape.text = ""
                text_frame = MagicMock()
                text_frame.word_wrap = True
                text_frame.auto_size = 1
                text_frame.vertical_anchor = 3
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.05)
                text_frame.margin_bottom = Inches(0.05)
                text_frame.paragraphs = [MagicMock()]
                text_frame.paragraphs[0].alignment = 0
                text_frame.paragraphs[0].font = MagicMock()
                text_frame.paragraphs[0].font.size = Pt(18)
                shape.text_frame = text_frame
                return shape

            mock_slide.shapes = [
                create_text_shape(2, "Title"),
                create_text_shape(4, "Content")
            ]

            # 为每个shape设置位置和尺寸属性
            for shape in mock_slide.shapes:
                shape.left = 914400  # 1英寸
                shape.top = 914400   # 1英寸
                shape.width = 914400  # 1英寸
                shape.height = 914400  # 1英寸

            strategy.create_slide("标题", "这是内容")

            assert strategy.slide_counter == 1

    def test_create_slide_empty_content(self, mock_presentation, config):
        """测试空内容"""
        with patch.object(mock_presentation.slides, 'add_slide') as mock_add:
            strategy = ContentSlideStrategy(mock_presentation, config)
            strategy.create_slide("标题", "")
            # 空内容不应该创建幻灯片
            mock_add.assert_not_called()

    def test_create_slide_whitespace_content(self, mock_presentation, config):
        """测试只有空白字符的内容"""
        with patch.object(mock_presentation.slides, 'add_slide') as mock_add:
            strategy = ContentSlideStrategy(mock_presentation, config)
            strategy.create_slide("标题", "   \n\n   ")
            mock_add.assert_not_called()

    def test_create_slide_long_content_split(self, mock_presentation, mock_slide, config):
        """测试长内容分块"""
        with patch.object(mock_presentation.slides, 'add_slide', return_value=mock_slide):
            strategy = ContentSlideStrategy(mock_presentation, config)

            # 创建多个文本形状
            def create_text_shape(shape_id, name):
                shape = MagicMock()
                shape.shape_id = shape_id
                shape.has_text_frame = True
                shape.name = name
                shape.is_placeholder = False
                shape.text = ""
                text_frame = MagicMock()
                text_frame.word_wrap = True
                text_frame.auto_size = 1
                text_frame.vertical_anchor = 3
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.05)
                text_frame.margin_bottom = Inches(0.05)
                text_frame.paragraphs = [MagicMock()]
                text_frame.paragraphs[0].alignment = 0
                text_frame.paragraphs[0].font = MagicMock()
                text_frame.paragraphs[0].font.size = Pt(18)
                shape.text_frame = text_frame
                return shape

            mock_slide.shapes = [
                create_text_shape(2, "Title"),
                create_text_shape(4, "Content")
            ]

            # 为每个shape设置位置和尺寸属性
            for shape in mock_slide.shapes:
                shape.left = 914400
                shape.top = 914400
                shape.width = 914400
                shape.height = 914400

            long_content = "这是一段很长的内容。" * 100
            strategy.create_slide("标题", long_content)

            # 应该创建多个幻灯片
            assert strategy.slide_counter >= 1

class TestTableOfContentsSlideStrategy:
    """测试TableOfContentsSlideStrategy"""

    def test_create_slide_with_3_items(self, mock_presentation, mock_slide, config):
        """测试3项目录"""
        mock_slide.shapes = [
            MagicMock(shape_id=1, has_text_frame=True, name="Title", is_placeholder=False, text=""),
            MagicMock(shape_id=2, has_text_frame=True, name="Item1", is_placeholder=False, text=""),
            MagicMock(shape_id=3, has_text_frame=True, name="Item2", is_placeholder=False, text=""),
            MagicMock(shape_id=4, has_text_frame=True, name="Item3", is_placeholder=False, text=""),
        ]

        # 为每个shape设置位置和尺寸属性
        for shape in mock_slide.shapes:
            shape.left = 914400
            shape.top = 914400
            shape.width = 914400
            shape.height = 914400

        with patch.object(mock_presentation.slides, 'add_slide', return_value=mock_slide):
            strategy = TableOfContentsSlideStrategy(mock_presentation, config)
            items = ["第一章", "第二章", "第三章"]
            strategy.create_slide(items)

            assert strategy.slide_counter == 1

    def test_create_slide_empty_items(self, mock_presentation, config):
        """测试空目录"""
        with patch.object(mock_presentation.slides, 'add_slide') as mock_add:
            strategy = TableOfContentsSlideStrategy(mock_presentation, config)
            strategy.create_slide([])
            mock_add.assert_not_called()

    def test_create_slide_many_items(self, mock_presentation, mock_slide, config):
        """测试多个目录项"""
        # 创建足够多的形状
        shapes = []
        for i in range(10):
            shape = MagicMock()
            shape.shape_id = i
            shape.has_text_frame = True
            shape.name = f"Shape{i}"
            shape.is_placeholder = False
            shape.text = ""
            text_frame = MagicMock()
            text_frame.word_wrap = True
            text_frame.auto_size = 1
            text_frame.vertical_anchor = 3
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            text_frame.paragraphs = [MagicMock()]
            text_frame.paragraphs[0].alignment = 0
            text_frame.paragraphs[0].font = MagicMock()
            text_frame.paragraphs[0].font.size = Pt(18)
            shape.text_frame = text_frame
            shapes.append(shape)

        mock_slide.shapes = shapes

        with patch.object(mock_presentation.slides, 'add_slide', return_value=mock_slide):
            strategy = TableOfContentsSlideStrategy(mock_presentation, config)
            items = [f"第{i}章" for i in range(1, 8)]
            strategy.create_slide(items)

            assert strategy.slide_counter == 1

class TestImageSlideStrategy:
    """测试ImageSlideStrategy"""

    def test_is_valid_image_url_valid_http(self):
        """测试有效的http URL"""
        config = SlideConfig()
        strategy = ImageSlideStrategy(None, config)
        assert strategy._is_valid_image_url("http://example.com/img.jpg") is True

    def test_is_valid_image_url_valid_https(self):
        """测试有效的https URL"""
        config = SlideConfig()
        strategy = ImageSlideStrategy(None, config)
        assert strategy._is_valid_image_url("https://example.com/img.png") is True

    def test_is_valid_image_url_invalid_protocol(self):
        """测试无效协议"""
        config = SlideConfig()
        strategy = ImageSlideStrategy(None, config)
        assert strategy._is_valid_image_url("ftp://example.com/img.jpg") is False

    def test_is_valid_image_url_empty(self):
        """测试空URL"""
        config = SlideConfig()
        strategy = ImageSlideStrategy(None, config)
        assert strategy._is_valid_image_url("") is False
        assert strategy._is_valid_image_url(None) is False

    @patch('utils.save_ppt.strategies.image_slide.requests.get')
    def test_download_image_success(self, mock_get):
        """测试图片下载成功"""
        mock_response = MagicMock()
        mock_response.content = b"fake image data"
        mock_response.raise_for_status = MagicMock()
        mock_get.return_value = mock_response

        config = SlideConfig()
        strategy = ImageSlideStrategy(None, config)

        result = strategy._download_image("http://example.com/img.jpg")

        assert result is not None
        mock_get.assert_called_once()

    @patch('utils.save_ppt.strategies.image_slide.requests.get')
    def test_download_image_failure(self, mock_get):
        """测试图片下载失败"""
        import requests
        mock_get.side_effect = requests.exceptions.RequestException("Network error")

        config = SlideConfig()
        strategy = ImageSlideStrategy(None, config)

        result = strategy._download_image("http://example.com/img.jpg")

        assert result is None

class TestSubSectionSlideStrategy:
    """测试SubSectionSlideStrategy"""

    def test_create_slide_with_3_items(self, mock_presentation, mock_slide, config):
        """测试3项子章节"""
        # 创建足够的文本形状
        shapes = []
        for i in range(10):
            shape = MagicMock()
            shape.shape_id = i
            shape.has_text_frame = True
            shape.name = f"Shape{i}"
            shape.is_placeholder = False
            shape.text = ""
            text_frame = MagicMock()
            text_frame.word_wrap = True
            text_frame.auto_size = 1
            text_frame.vertical_anchor = 3
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            text_frame.paragraphs = [MagicMock()]
            text_frame.paragraphs[0].alignment = 0
            text_frame.paragraphs[0].font = MagicMock()
            text_frame.paragraphs[0].font.size = Pt(18)
            shape.text_frame = text_frame
            shapes.append(shape)

        mock_slide.shapes = shapes

        with patch.object(mock_presentation.slides, 'add_slide', return_value=mock_slide):
            strategy = SubSectionSlideStrategy(mock_presentation, config)
            content = [
                {"summary": "要点1", "detail": "详情1"},
                {"summary": "要点2", "detail": "详情2"},
                {"summary": "要点3", "detail": "详情3"}
            ]
            strategy.create_slide("章节标题", content)

            assert strategy.slide_counter == 1

    def test_create_slide_empty_content(self, mock_presentation, config):
        """测试空内容"""
        with patch.object(mock_presentation.slides, 'add_slide') as mock_add:
            strategy = SubSectionSlideStrategy(mock_presentation, config)
            strategy.create_slide("标题", [])
            mock_add.assert_not_called()

    def test_create_slide_unsupported_count(self, mock_presentation, config):
        """测试不支持的数量"""
        with patch.object(mock_presentation.slides, 'add_slide') as mock_add:
            strategy = SubSectionSlideStrategy(mock_presentation, config)
            content = [{"summary": f"要点{i}", "detail": f"详情{i}"} for i in range(10)]
            strategy.create_slide("标题", content)
            # 超过5项不应该创建幻灯片
            mock_add.assert_not_called()

class TestReferencesSlideStrategy:
    """测试ReferencesSlideStrategy"""

    def test_process_reference_text_basic(self, mock_presentation, config):
        """测试基本参考文献处理"""
        strategy = ReferencesSlideStrategy(mock_presentation, config)

        ref = "Author. Title. Journal, 2023."
        result = strategy._process_reference_text(ref)

        assert "Author" in result
        assert "Title" in result

    def test_process_reference_text_short_line(self, mock_presentation, config):
        """测试短行合并"""
        strategy = ReferencesSlideStrategy(mock_presentation, config)

        ref = "Vol. 1\nNo. 2\nAuthor. Title. Journal, 2023."
        result = strategy._process_reference_text(ref)

        # 短行应该被合并
        assert "Vol. 1" in result

    def test_create_slide_with_references(self, mock_presentation, mock_slide, config):
        """测试创建参考文献页"""
        # 创建足够的形状
        shapes = []
        for i in range(15):
            shape = MagicMock()
            shape.shape_id = i
            shape.has_text_frame = True
            shape.name = f"Shape{i}"
            shape.is_placeholder = False
            shape.text = ""
            text_frame = MagicMock()
            text_frame.word_wrap = True
            text_frame.auto_size = 1
            text_frame.vertical_anchor = 3
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            text_frame.paragraphs = [MagicMock()]
            text_frame.paragraphs[0].alignment = 0
            text_frame.paragraphs[0].font = MagicMock()
            text_frame.paragraphs[0].font.size = Pt(18)
            shape.text_frame = text_frame
            shapes.append(shape)

        mock_slide.shapes = shapes

        with patch.object(mock_presentation.slides, 'add_slide', return_value=mock_slide):
            strategy = ReferencesSlideStrategy(mock_presentation, config)
            refs = [
                "Reference 1: Author, Title, Journal, 2023",
                "Reference 2: Author2, Title2, Journal2, 2022"
            ]
            strategy.create_slide(refs)

            assert strategy.slide_counter >= 1

    def test_create_slide_empty_references(self, mock_presentation, config):
        """测试空参考文献"""
        with patch.object(mock_presentation.slides, 'add_slide') as mock_add:
            strategy = ReferencesSlideStrategy(mock_presentation, config)
            strategy.create_slide([])
            mock_add.assert_not_called()

class TestEndSlideStrategy:
    """测试EndSlideStrategy"""

    def test_create_slide(self, mock_presentation, mock_slide, config):
        """测试创建结束页"""
        with patch.object(mock_presentation.slides, 'add_slide', return_value=mock_slide):
            strategy = EndSlideStrategy(mock_presentation, config)
            strategy.create_slide()

            assert strategy.slide_counter == 1
