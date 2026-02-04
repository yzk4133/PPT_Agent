"""
幻灯片生成策略抽象基类

包含所有策略共享的方法和抽象接口。
"""

import logging
from typing import Optional
from abc import ABC, abstractmethod
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR

from ..config import SlideConfig
from ..text_processor import TextProcessor

logger = logging.getLogger(__name__)

class SlideStrategy(ABC):
    """幻灯片生成策略的抽象基类"""

    def __init__(self, presentation: Presentation, config: SlideConfig):
        self.presentation = presentation
        self.config = config
        self.text_processor = TextProcessor()
        self.slide_counter = 0  # 添加页面计数器

    @abstractmethod
    def create_slide(self, *args, **kwargs):
        """创建幻灯片的抽象方法"""
        pass

    def _log_slide_shapes(self, slide, slide_type: str):
        """记录幻灯片上所有形状的详细信息"""
        logger.info(f"\n{'='*80}")
        logger.info(f"正在分析第 {self.slide_counter} 页 - 类型: {slide_type}")
        logger.info(f"{'='*80}")

        logger.info(f"幻灯片上共有 {len(slide.shapes)} 个形状:")
        for idx, shape in enumerate(slide.shapes):
            shape_info = []
            shape_info.append(f"  形状 #{idx + 1}:")
            shape_info.append(f"    - 形状ID: {shape.shape_id}")
            shape_info.append(f"    - 名称: {shape.name}")
            shape_info.append(f"    - 类型: {shape.shape_type}")

            # 检查是否有位置和尺寸属性，并且是否是数字类型（非Mock）
            try:
                if hasattr(shape, 'left') and hasattr(shape, 'top'):
                    left_val = shape.left
                    top_val = shape.top
                    # 检查是否可以转换为数字（排除Mock对象）
                    if isinstance(left_val, (int, float)) and isinstance(top_val, (int, float)):
                        shape_info.append(f"    - 位置: 左={left_val/914400:.2f}英寸, 上={top_val/914400:.2f}英寸")

                if hasattr(shape, 'width') and hasattr(shape, 'height'):
                    width_val = shape.width
                    height_val = shape.height
                    if isinstance(width_val, (int, float)) and isinstance(height_val, (int, float)):
                        shape_info.append(f"    - 尺寸: 宽={width_val/914400:.2f}英寸, 高={height_val/914400:.2f}英寸")
                        # 添加形状面积信息，帮助比较
                        area = (width_val/914400) * (height_val/914400)
                        shape_info.append(f"    - 面积: {area:.2f}平方英寸")
            except (TypeError, AttributeError):
                # 忽略Mock对象或属性访问错误
                pass

            shape_info.append(f"    - 是占位符: {shape.is_placeholder}")
            shape_info.append(f"    - 有文本框: {shape.has_text_frame}")

            if shape.has_text_frame:
                text_frame = shape.text_frame
                shape_info.append(f"    - 文本框设置:")
                shape_info.append(f"      - 自动调整: {text_frame.auto_size}")
                shape_info.append(f"      - 垂直锚点: {text_frame.vertical_anchor}")
                shape_info.append(f"      - 自动换行: {text_frame.word_wrap}")
                shape_info.append(f"      - 边距: 左={text_frame.margin_left}, 右={text_frame.margin_right}")

                if shape.text:
                    text_preview = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                    shape_info.append(f"    - 文本预览: '{text_preview}'")

            logger.info('\n'.join(shape_info))
        logger.info(f"{'='*80}\n")

    def _get_slide_layout(self, layout_key: str):
        """获取幻灯片布局"""
        layout_id = self.config.SLIDE_LAYOUTS.get(layout_key)
        if layout_id is None or layout_id >= len(self.presentation.slide_layouts):
            logger.warning(f"布局 '{layout_key}' (ID: {layout_id}) 未找到，使用默认布局")
            logger.warning(f"可用布局数量: {len(self.presentation.slide_layouts)}")
            return self.presentation.slide_layouts[0]

        print(f"使用布局: '{layout_key}' (ID: {layout_id})")
        return self.presentation.slide_layouts[layout_id]

    def _add_text_with_auto_fit(
        self,
        slide,
        shape_id: int,
        text: str,
        font_type: str = "content",
        max_chars: Optional[int] = None,
        is_title_page: bool = False
    ):
        """添加文本并自动调整字体大小，同时处理标题背景。"""
        logger.debug(f"尝试向形状ID {shape_id} 添加文本: '{text[:30]}...' (类型: {font_type})")

        shape_found = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_id == shape_id:
                shape_found = True
                clean_text = self.text_processor.remove_html_tags(text)

                if max_chars:
                    original_length = len(clean_text)
                    clean_text = self.text_processor.truncate_text(clean_text, max_chars)
                    if original_length > max_chars:
                        logger.debug(f"文本被截断: {original_length} -> {len(clean_text)} 字符")

                shape.text = clean_text
                text_frame = shape.text_frame

                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.05)
                text_frame.margin_bottom = Inches(0.05)

                # 1. 非首页的标题
                if font_type == "title" and not is_title_page:
                    text_frame.word_wrap = False
                    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                    font_size = self.config.FONT_SIZES["title"]["default"]

                    for p in text_frame.paragraphs:
                        # p.alignment = PP_ALIGN.LEFT # 标题居左
                        p.alignment = 1  # PP_ALIGN.CENTER
                        p.font.size = Pt(font_size)

                    self._adjust_title_background_shape(slide, shape, clean_text, font_size)

                # 2. 首页的标题
                elif font_type == "title" and is_title_page:
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    for p in text_frame.paragraphs:
                        p.alignment = 1  # PP_ALIGN.CENTER

                # 3. 其他所有内容文本
                else:
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                print(f"✓ 成功添加文本到形状 {shape_id} (名称: {shape.name})，文本长度: {len(clean_text)} 字符")
                return

        if not shape_found:
            logger.error(f"✗ 未找到形状ID {shape_id} - 请检查模板中的形状ID设置")
            # 列出所有可用的形状ID
            available_ids = [s.shape_id for s in slide.shapes if s.has_text_frame]
            logger.error(f"  可用的文本形状ID: {available_ids}")

    def _adjust_title_background_shape(self, slide, text_shape, text: str, font_size: int):
        """调整标题背景形状"""

        try:
            layout = slide.slide_layout
            layout_index = self.presentation.slide_layouts.index(layout)
            layout_name = layout.name or "未命名"

            print(f"使用的母版布局: 索引={layout_index}, 名称='{layout_name}'")
            logger.debug(f"母版布局详细信息: {layout}")
            bg_shape = None
            if layout_index == 16:
                target_shape_name = "Text Placeholder 4"
            else:
                target_shape_name = "Text Placeholder 1"

            logger.debug(f"查找背景形状: '{target_shape_name}'")

            for shape in slide.shapes:
                logger.debug(f"  检查形状: {shape.name} (ID: {shape.shape_id})")
                if shape.name == target_shape_name:
                    bg_shape = shape
                    print(f"✓ 找到背景形状: '{target_shape_name}'")
                    break

            if not bg_shape:
                logger.warning(f"✗ 未找到名为 '{target_shape_name}' 的背景形状")
                return

            # 记录原始尺寸
            original_bg_left = bg_shape.left
            original_bg_top = bg_shape.top
            original_bg_height = bg_shape.height
            original_offset_x = text_shape.left - original_bg_left

            logger.debug(f"  背景形状原始位置: 左={original_bg_left/914400:.2f}英寸, 上={original_bg_top/914400:.2f}英寸")
            logger.debug(f"  背景形状原始尺寸: 宽={bg_shape.width/914400:.2f}英寸, 高={original_bg_height/914400:.2f}英寸")

            if bg_shape.has_text_frame:
                bg_shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE

            # 计算新宽度
            calculated_text_width = self._calculate_text_width(text, font_size)
            padding = Inches(self.config.BACKGROUND_SHAPE_CONFIG["padding"])
            text_margin_padding = text_shape.text_frame.margin_left + text_shape.text_frame.margin_right
            new_width = calculated_text_width + text_margin_padding + padding

            min_width = Inches(self.config.BACKGROUND_SHAPE_CONFIG["min_width"])
            new_width = max(new_width, min_width)

            # 应用新尺寸
            bg_shape.width = int(new_width)
            bg_shape.left = text_shape.left - original_offset_x
            bg_shape.top = original_bg_top
            bg_shape.height = original_bg_height

            # 边界检查
            slide_width = self.presentation.slide_width
            if bg_shape.left < Inches(0.2):
                bg_shape.left = int(Inches(0.2))
            if bg_shape.left + bg_shape.width > slide_width - Inches(0.2):
                bg_shape.left = int(slide_width - bg_shape.width - Inches(0.2))

            print(f"✓ 背景形状调整完成: 新宽度={new_width/914400:.2f}英寸")

        except Exception as e:
            logger.error(f"✗ 调整背景图形失败: {e}", exc_info=True)

    def _calculate_text_width(self, text: str, font_size: int) -> int:
        """精确计算文本的实际宽度（返回EMU单位）"""
        try:
            # 字符分类统计
            chinese_chars = 0
            english_chars = 0
            numbers = 0
            spaces = 0
            other_chars = 0

            for char in text:
                if '\u4e00' <= char <= '\u9fff':  # 中文字符
                    chinese_chars += 1
                elif char in '，。！？；：""''（）【】《》、':  # 中文标点
                    chinese_chars += 1
                elif char.isalpha():  # 英文字母
                    english_chars += 1
                elif char.isdigit():  # 数字
                    numbers += 1
                elif char == ' ':  # 空格
                    spaces += 1
                else:  # 其他字符
                    other_chars += 1

            # 不同字符类型的宽度系数
            width_factors = {
                'chinese': 0.55,
                'english': 0.55,
                'numbers': 0.6,
                'spaces': 0.3,
                'others': 0.7
            }

            # 计算总宽度
            total_width_pt = (
                chinese_chars * font_size * width_factors['chinese'] +
                english_chars * font_size * width_factors['english'] +
                numbers * font_size * width_factors['numbers'] +
                spaces * font_size * width_factors['spaces'] +
                other_chars * font_size * width_factors['others']
            )

            # 添加字符间距的影响
            total_width_pt *= 1.05

            # 转换为EMU
            width_in_emu = int(total_width_pt * 914400 / 72)

            logger.debug(f"文本宽度计算: 中文={chinese_chars}, 英文={english_chars}, "
                        f"字体={font_size}pt, 宽度={width_in_emu/914400:.2f}英寸")

            return width_in_emu

        except Exception as e:
            logger.error(f"计算文本宽度失败: {e}")
            return Inches(4)

    def _fill_empty_placeholders(self, slide):
        """移除空占位符"""
        removed_count = 0
        for shape in list(slide.shapes):
            if shape.has_text_frame and not shape.text.strip() and shape.is_placeholder:
                sp = shape._element
                sp.getparent().remove(sp)
                removed_count += 1

        if removed_count > 0:
            logger.debug(f"移除了 {removed_count} 个空占位符")
