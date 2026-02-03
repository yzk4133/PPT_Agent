from pptx import Presentation

new = Presentation('ppt_template_0717.pptx')

for idx, layout in enumerate(new.slide_layouts):
    print(f"\n--- 母版布局 {idx} [{layout.name}] ---")
    for shape in layout.shapes:
        print(f"ID: {shape.shape_id} | 名称: {shape.name} | 类型: {shape.shape_type}")

old = Presentation('ppt_template_0717.pptx')

for idx, layout in enumerate(old.slide_layouts):
    print(f"\n--- 母版布局 {idx} [{layout.name}] ---")
    for shape in layout.shapes:
        print(f"ID: {shape.shape_id} | 名称: {shape.name} | 类型: {shape.shape_type}")