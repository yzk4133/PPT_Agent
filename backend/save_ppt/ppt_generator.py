import os
import json
import logging
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from abc import ABC, abstractmethod
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
import random
import re
import requests
from io import BytesIO
from PIL import Image
import datetime

# 配置更详细的日志格式
logging.basicConfig(
    level=logging.DEBUG,  # 改为DEBUG级别
    format='%(asctime)s - %(name)s - %(levelname)s - %(funcName)s:%(lineno)d - %(message)s'
)
logger = logging.getLogger(__name__)

# ==================== 配置类 ====================
@dataclass
class SlideConfig:
    """幻灯片配置数据类"""
    SLIDE_LAYOUTS = {
        "TITLE_PAGE": 0,
        "CONTENT_TITLE_AND_TEXT": 1,
        "TABLE_OF_CONTENTS_GENERIC": 22,
        "TABLE_OF_CONTENTS_4_ITEMS_A": 23,
        "TABLE_OF_CONTENTS_4_ITEMS_B": 24,
        "TABLE_OF_CONTENTS_5_ITEMS_A": 25,
        "TABLE_OF_CONTENTS_5_ITEMS_B": 26,
        "TABLE_OF_CONTENTS_3_ITEMS": 30,
        "INFO_TITLE_AND_TEXT": 15,
        "TEXT_ONLY_SMALL_TITLE": 12,
        "IMAGE_TITLE_AND_DESCRIPTION_WIDE": 9,
        "IMAGE_TITLE_AND_DESCRIPTION_TALL": 7,
        "IMAGE_ONLY": 14,
        "REFERENCES_PAGE": 10,
        "END_PAGE": 11,
        "SUBCHAPTER_2_ITEMS": 15,
        "SUBCHAPTER_3_ITEMS": 16,
        "SUBCHAPTER_4_ITEMS": 17,
        "SUBCHAPTER_5_ITEMS": 10,
    }
    
    
    SHAPE_IDS = {
        # 标题页
        "TITLE_PAGE_TITLE": 2,
        "TITLE_PAGE_DATE": 9,
        
        # 目录页
        "TOC_TITLE": 1,
        "TOC_ITEMS": {1: 2, 2: 3, 3: 4, 4: 5, 5: 6, 6: 7},
        
        # 内容页
        "CONTENT_TITLE": 2,
        "CONTENT_TEXT": 3,
        "CONTENT_TEXT_LAYOUT1_4": 4,
        
        # 图片页
        "IMAGE_PLACEHOLDER": 3,
        "IMAGE_TITLE": 2,
        "IMAGE_DESCRIPTION": 4,
        "IMAGE_DIO_LAYOUT7": 6,
        "IMAGE_DIO_LAYOUT9": 5,
        
        # 参考文献页
        "REFERENCES_TITLE": 2,
        "REFERENCES": [
            {"num": 3, "text": 8},
            {"num": 4, "text": 9},
            {"num": 5, "text": 10},
            {"num": 6, "text": 11},
            {"num": 7, "text": 12},
        ],
        
        # 子章节
        "SUBCHAPTER_TITLE": 5,
        "SUBCHAPTER_ITEMS": {
            1: [2],
            2: [2, 3],
            3: [2, 3, 4],
            4: [2, 3, 4, 5],
            5: [(3, 8), (4, 9), (5, 10), (6, 11), (7, 12)],
        }
    }
    
    # 文本处理常量
    MAX_CHARS_PER_TEXT_CHUNK = 800
    MAX_REFERENCES_PER_SLIDE = 5
    MAX_TOTAL_REFERENCES = 10
    
    # 图片URL前缀
    VALID_IMAGE_URL_PREFIXES = [
        "http",
        "https",
    ]
    
    # 字体大小配置
    FONT_SIZES = {
        "title": {"min": 18, "max": 44, "default": 26},
        "content": {"min": 12, "max": 24, "default": 18},
        "small": {"min": 10, "max": 16, "default": 14},
    }
    
    # 背景图形调整配置
    BACKGROUND_SHAPE_CONFIG = {
        "min_width": 2.5,          # 最小宽度（英寸）
        "padding": 0.8,            # 总padding（英寸）
        "vertical_tolerance": 0.5,  # 垂直位置容差（英寸）
        "height_tolerance": 0.3,    # 高度容差（英寸）
        "alignment": "center"       # 对齐方式：center, left, right
    }

# ==================== 文本处理器 ====================
class TextProcessor:
    """文本处理工具类"""
    
    @staticmethod
    def remove_html_tags(text: str) -> str:
        """去除HTML标签"""
        if not isinstance(text, str):
            return ""
        clean = re.sub(r'<.*?>', '', text)
        return clean.strip()
    
    @staticmethod
    def calculate_optimal_font_size(
        text: str, 
        shape,
        font_type: str = "content"
    ) -> int:
        """根据文本长度和文本框大小计算最佳字体大小"""
        if not text or not shape:
            return SlideConfig.FONT_SIZES[font_type]["default"]
        
        try:
            width = shape.width
            height = shape.height
        except:
            return SlideConfig.FONT_SIZES[font_type]["default"]
        
        char_count = len(text)
        estimated_chars_per_line = int(width / Pt(7))
        estimated_lines = max(1, char_count / max(1, estimated_chars_per_line))
        font_config = SlideConfig.FONT_SIZES[font_type]
        
        if estimated_lines > 20:
            return font_config["min"]
        elif estimated_lines > 10:
            return int(font_config["min"] + (font_config["default"] - font_config["min"]) * 0.5)
        elif estimated_lines > 5:
            return font_config["default"]
        else:
            return min(font_config["max"], font_config["default"] + 4)
    
    @staticmethod
    def truncate_text(text: str, max_chars: int, suffix: str = "...") -> str:
        """截断过长的文本"""
        if len(text) <= max_chars:
            return text
        return text[:max_chars - len(suffix)] + suffix
    
    @staticmethod
    def split_text_into_chunks(
        content: str, 
        max_chars: int = SlideConfig.MAX_CHARS_PER_TEXT_CHUNK
    ) -> List[str]:
        """将长文本分割成块"""
        if not content or not content.strip():
            return []
        
        sentences = re.split(r'(?<=[.!?。？！]) +', content)
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            if len(current_chunk) + len(sentence) + 1 > max_chars and current_chunk:
                chunks.append(current_chunk.strip())
                current_chunk = sentence + " "
            else:
                current_chunk += sentence + " "
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks

# ==================== 幻灯片策略接口 ====================
class SlideStrategy(ABC):
    """幻灯片生成策略的抽象基类"""
    
    def __init__(self, presentation: Presentation, config: SlideConfig):
        self.presentation = presentation
        self.config = config
        self.text_processor = TextProcessor()
        self.slide_counter = 0  # 添加页面计数器
    
    @abstractmethod
    def create_slide(self, *args, **kwargs):
        """创建幻灯片的抽象方法"""
        pass
    
    def _log_slide_shapes(self, slide, slide_type: str):
        """记录幻灯片上所有形状的详细信息"""
        logger.info(f"\n{'='*80}")
        logger.info(f"正在分析第 {self.slide_counter} 页 - 类型: {slide_type}")
        logger.info(f"{'='*80}")
        
        logger.info(f"幻灯片上共有 {len(slide.shapes)} 个形状:")
        for idx, shape in enumerate(slide.shapes):
            shape_info = []
            shape_info.append(f"  形状 #{idx + 1}:")
            shape_info.append(f"    - 形状ID: {shape.shape_id}")
            shape_info.append(f"    - 名称: {shape.name}")
            shape_info.append(f"    - 类型: {shape.shape_type}")
            
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                shape_info.append(f"    - 位置: 左={shape.left/914400:.2f}英寸, 上={shape.top/914400:.2f}英寸")
            
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                shape_info.append(f"    - 尺寸: 宽={shape.width/914400:.2f}英寸, 高={shape.height/914400:.2f}英寸")
                # 添加形状面积信息，帮助比较
                area = (shape.width/914400) * (shape.height/914400)
                shape_info.append(f"    - 面积: {area:.2f}平方英寸")
            
            shape_info.append(f"    - 是占位符: {shape.is_placeholder}")
            shape_info.append(f"    - 有文本框: {shape.has_text_frame}")
            
            if shape.has_text_frame:
                text_frame = shape.text_frame
                shape_info.append(f"    - 文本框设置:")
                shape_info.append(f"      - 自动调整: {text_frame.auto_size}")
                shape_info.append(f"      - 垂直锚点: {text_frame.vertical_anchor}")
                shape_info.append(f"      - 自动换行: {text_frame.word_wrap}")
                shape_info.append(f"      - 边距: 左={text_frame.margin_left}, 右={text_frame.margin_right}")
                
                if shape.text:
                    text_preview = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                    shape_info.append(f"    - 文本预览: '{text_preview}'")
            
            logger.info('\n'.join(shape_info))
        logger.info(f"{'='*80}\n")


    def _get_slide_layout(self, layout_key: str):
        """获取幻灯片布局"""
        layout_id = self.config.SLIDE_LAYOUTS.get(layout_key)
        if layout_id is None or layout_id >= len(self.presentation.slide_layouts):
            logger.warning(f"布局 '{layout_key}' (ID: {layout_id}) 未找到，使用默认布局")
            logger.warning(f"可用布局数量: {len(self.presentation.slide_layouts)}")
            return self.presentation.slide_layouts[0]
        
        print(f"使用布局: '{layout_key}' (ID: {layout_id})")
        return self.presentation.slide_layouts[layout_id]
    
    def _add_text_with_auto_fit(
        self,
        slide,
        shape_id: int,
        text: str,
        font_type: str = "content",
        max_chars: Optional[int] = None,
        is_title_page: bool = False
    ):
        """添加文本并自动调整字体大小，同时处理标题背景。"""
        logger.debug(f"尝试向形状ID {shape_id} 添加文本: '{text[:30]}...' (类型: {font_type})")
        
        shape_found = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_id == shape_id:
                shape_found = True
                clean_text = self.text_processor.remove_html_tags(text)

                if max_chars:
                    original_length = len(clean_text)
                    clean_text = self.text_processor.truncate_text(clean_text, max_chars)
                    if original_length > max_chars:
                        logger.debug(f"文本被截断: {original_length} -> {len(clean_text)} 字符")

                shape.text = clean_text
                text_frame = shape.text_frame

                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.05)
                text_frame.margin_bottom = Inches(0.05)

                # 1. 非首页的标题
                if font_type == "title" and not is_title_page:
                    text_frame.word_wrap = False
                    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    font_size = self.config.FONT_SIZES["title"]["default"]
                    
                    for p in text_frame.paragraphs:
                        # p.alignment = PP_ALIGN.LEFT # 标题居左
                        p.alignment = PP_ALIGN.CENTER 
                        p.font.size = Pt(font_size)

                    self._adjust_title_background_shape(slide, shape, clean_text, font_size)

                # 2. 首页的标题
                elif font_type == "title" and is_title_page:
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    for p in text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER

                # 3. 其他所有内容文本
                else:
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                print(f"✓ 成功添加文本到形状 {shape_id} (名称: {shape.name})，文本长度: {len(clean_text)} 字符")
                return

        if not shape_found:
            logger.error(f"✗ 未找到形状ID {shape_id} - 请检查模板中的形状ID设置")
            # 列出所有可用的形状ID
            available_ids = [s.shape_id for s in slide.shapes if s.has_text_frame]
            logger.error(f"  可用的文本形状ID: {available_ids}")

    def _adjust_title_background_shape(self, slide, text_shape, text: str, font_size: int):
        """调整标题背景形状"""
        
        try:
            layout = slide.slide_layout
            layout_index = self.presentation.slide_layouts.index(layout)
            layout_name = layout.name or "未命名"
            
            print(f"使用的母版布局: 索引={layout_index}, 名称='{layout_name}'")
            logger.debug(f"母版布局详细信息: {layout}")
            bg_shape = None
            if layout_index == 16:
                target_shape_name = "Text Placeholder 4"
            else:
                target_shape_name = "Text Placeholder 1" 
            
            logger.debug(f"查找背景形状: '{target_shape_name}'")
            
            for shape in slide.shapes:
                logger.debug(f"  检查形状: {shape.name} (ID: {shape.shape_id})")
                if shape.name == target_shape_name:
                    bg_shape = shape
                    print(f"✓ 找到背景形状: '{target_shape_name}'")
                    break
            
            if not bg_shape:
                logger.warning(f"✗ 未找到名为 '{target_shape_name}' 的背景形状")
                return

            # 记录原始尺寸
            original_bg_left = bg_shape.left
            original_bg_top = bg_shape.top
            original_bg_height = bg_shape.height
            original_offset_x = text_shape.left - original_bg_left
            
            logger.debug(f"  背景形状原始位置: 左={original_bg_left/914400:.2f}英寸, 上={original_bg_top/914400:.2f}英寸")
            logger.debug(f"  背景形状原始尺寸: 宽={bg_shape.width/914400:.2f}英寸, 高={original_bg_height/914400:.2f}英寸")
            
            if bg_shape.has_text_frame:
                bg_shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE

            # 计算新宽度
            calculated_text_width = self._calculate_text_width(text, font_size)
            padding = Inches(self.config.BACKGROUND_SHAPE_CONFIG["padding"])
            text_margin_padding = text_shape.text_frame.margin_left + text_shape.text_frame.margin_right
            new_width = calculated_text_width + text_margin_padding + padding
            
            min_width = Inches(self.config.BACKGROUND_SHAPE_CONFIG["min_width"])
            new_width = max(new_width, min_width)
            
            # 应用新尺寸
            bg_shape.width = int(new_width)
            bg_shape.left = text_shape.left - original_offset_x
            bg_shape.top = original_bg_top
            bg_shape.height = original_bg_height

            # 边界检查
            slide_width = self.presentation.slide_width
            if bg_shape.left < Inches(0.2):
                bg_shape.left = int(Inches(0.2))
            if bg_shape.left + bg_shape.width > slide_width - Inches(0.2):
                bg_shape.left = int(slide_width - bg_shape.width - Inches(0.2))

            print(f"✓ 背景形状调整完成: 新宽度={new_width/914400:.2f}英寸")

        except Exception as e:
            logger.error(f"✗ 调整背景图形失败: {e}", exc_info=True)

    def _calculate_text_width(self, text: str, font_size: int) -> int:
        """精确计算文本的实际宽度（返回EMU单位）"""
        try:
            # 字符分类统计
            chinese_chars = 0
            english_chars = 0
            numbers = 0
            spaces = 0
            other_chars = 0
            
            for char in text:
                if '\u4e00' <= char <= '\u9fff':  # 中文字符
                    chinese_chars += 1
                elif char in '，。！？；：""''（）【】《》、':  # 中文标点
                    chinese_chars += 1
                elif char.isalpha():  # 英文字母
                    english_chars += 1
                elif char.isdigit():  # 数字
                    numbers += 1
                elif char == ' ':  # 空格
                    spaces += 1
                else:  # 其他字符
                    other_chars += 1
            
            # 不同字符类型的宽度系数
            width_factors = {
                'chinese': 0.55,
                'english': 0.55,
                'numbers': 0.6,
                'spaces': 0.3,
                'others': 0.7
            }
            
            # 计算总宽度
            total_width_pt = (
                chinese_chars * font_size * width_factors['chinese'] +
                english_chars * font_size * width_factors['english'] +
                numbers * font_size * width_factors['numbers'] +
                spaces * font_size * width_factors['spaces'] +
                other_chars * font_size * width_factors['others']
            )
            
            # 添加字符间距的影响
            total_width_pt *= 1.05
            
            # 转换为EMU
            width_in_emu = int(total_width_pt * 914400 / 72)
            
            logger.debug(f"文本宽度计算: 中文={chinese_chars}, 英文={english_chars}, "
                        f"字体={font_size}pt, 宽度={width_in_emu/914400:.2f}英寸")
            
            return width_in_emu
            
        except Exception as e:
            logger.error(f"计算文本宽度失败: {e}")
            return Inches(4)

    def _fill_empty_placeholders(self, slide):
        """移除空占位符"""
        removed_count = 0
        for shape in list(slide.shapes):
            if shape.has_text_frame and not shape.text.strip() and shape.is_placeholder:
                sp = shape._element
                sp.getparent().remove(sp)
                removed_count += 1
        
        if removed_count > 0:
            logger.debug(f"移除了 {removed_count} 个空占位符")

# ==================== 具体的幻灯片策略实现 ====================
class TitleSlideStrategy(SlideStrategy):
    """标题页生成策略"""
    
    def create_slide(self, title: str):
        self.slide_counter += 1
        print(f"\n{'#'*80}")
        print(f"# 创建第 {self.slide_counter} 页: 标题页")
        print(f"# 标题: {title}")
        print(f"{'#'*80}")
        
        slide_layout = self._get_slide_layout("TITLE_PAGE")
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 记录所有形状信息
        self._log_slide_shapes(slide, "标题页")
        
        # 添加标题
        self._add_text_with_auto_fit(
            slide, 
            self.config.SHAPE_IDS["TITLE_PAGE_TITLE"], 
            title,
            font_type="title",
            max_chars=None,
            is_title_page=True
        )
        
        # 添加日期
        current_date = datetime.datetime.now().strftime("%Y年%m月%d日")
        self._add_text_with_auto_fit(
            slide,
            self.config.SHAPE_IDS["TITLE_PAGE_DATE"],
            current_date,
            font_type="small"
        )
        
        self._fill_empty_placeholders(slide)
        print(f"✓ 标题页创建完成")

class ContentSlideStrategy(SlideStrategy):
    """内容页生成策略"""
    
    def create_slide(self, title: str, content: str):
        if not content or not content.strip():
            logger.debug(f"跳过空内容页: '{title}'")
            return
        
        content_chunks = self.text_processor.split_text_into_chunks(content)
        print(f"内容被分为 {len(content_chunks)} 个块")
        
        for idx, chunk in enumerate(content_chunks):
            self.slide_counter += 1
            slide_title = title if idx == 0 else f"{title} (续 {idx + 1})"
            
            print(f"\n{'#'*80}")
            print(f"# 创建第 {self.slide_counter} 页: 内容页")
            print(f"# 标题: {slide_title}")
            print(f"# 内容长度: {len(chunk)} 字符")
            print(f"{'#'*80}")
            
            # 根据内容长度选择布局
            if len(chunk) < 150:
                layout_key = "TEXT_ONLY_SMALL_TITLE"
                text_shape_id = self.config.SHAPE_IDS["CONTENT_TEXT"]
                print("使用小标题布局（内容较短）")
            else:
                layout_key = "CONTENT_TITLE_AND_TEXT"
                text_shape_id = self.config.SHAPE_IDS["CONTENT_TEXT_LAYOUT1_4"]
                print("使用标准内容布局")
            
            slide_layout = self._get_slide_layout(layout_key)
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # 记录所有形状信息
            self._log_slide_shapes(slide, "内容页")
            
            # 添加标题
            self._add_text_with_auto_fit(
                slide,
                self.config.SHAPE_IDS["CONTENT_TITLE"],
                slide_title,
                font_type="title",
                max_chars=80,
                is_title_page=False
            )
            
            # 添加内容
            self._add_text_with_auto_fit(
                slide,
                text_shape_id,
                chunk,
                font_type="content"
            )
            
            self._fill_empty_placeholders(slide)
            print(f"✓ 内容页创建完成")

class TableOfContentsSlideStrategy(SlideStrategy):
    """目录页生成策略"""
    
    def create_slide(self, toc_items: List[str]):
        if not toc_items:
            print("无目录项，跳过目录页")
            return
        
        self.slide_counter += 1
        print(f"\n{'#'*80}")
        print(f"# 创建第 {self.slide_counter} 页: 目录页")
        print(f"# 目录项数: {len(toc_items)}")
        print(f"{'#'*80}")
        
        # 根据项目数选择布局
        layout_map = {
            3: "TABLE_OF_CONTENTS_3_ITEMS",
            4: random.choice(["TABLE_OF_CONTENTS_4_ITEMS_A", "TABLE_OF_CONTENTS_4_ITEMS_B"]),
            5: random.choice(["TABLE_OF_CONTENTS_5_ITEMS_A", "TABLE_OF_CONTENTS_5_ITEMS_B"]),
        }
        
        layout_key = layout_map.get(len(toc_items), "TABLE_OF_CONTENTS_GENERIC")
        print(f"根据 {len(toc_items)} 个目录项，选择布局: {layout_key}")
        
        slide_layout = self._get_slide_layout(layout_key)
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 记录所有形状信息
        self._log_slide_shapes(slide, "目录页")
        
        # 添加标题
        self._add_text_with_auto_fit(
            slide,
            self.config.SHAPE_IDS["TOC_TITLE"],
            "目录",
            font_type="title",
            is_title_page=False
        )
        
        # 添加目录项
        toc_shape_ids = self.config.SHAPE_IDS["TOC_ITEMS"]
        for i, item in enumerate(toc_items[:6]):  # 最多6项
            if i + 1 in toc_shape_ids:
                print(f"添加目录项 {i+1}: {item}")
                self._add_text_with_auto_fit(
                    slide,
                    toc_shape_ids[i + 1],
                    item,
                    font_type="content",
                    max_chars=80
                )
        
        self._fill_empty_placeholders(slide)
        print(f"✓ 目录页创建完成")

class ImageSlideStrategy(SlideStrategy):
    """图片页生成策略"""
    
    def _is_valid_image_url(self, url: str) -> bool:
        """检查图片URL是否有效"""
        if not isinstance(url, str):
            return False
        is_valid = any(url.startswith(prefix) for prefix in self.config.VALID_IMAGE_URL_PREFIXES)
        logger.debug(f"图片URL验证: {url} -> {'有效' if is_valid else '无效'}")
        return is_valid
    
    def _download_image(self, image_url: str) -> Optional[BytesIO]:
        """下载图片"""
        if not self._is_valid_image_url(image_url):
            logger.warning(f"无效的图片URL: {image_url}")
            return None
        
        try:
            print(f"开始下载图片: {image_url}")
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            print(f"✓ 成功下载图片，大小: {len(response.content)} 字节")
            return BytesIO(response.content)
        except requests.exceptions.RequestException as e:
            logger.error(f"✗ 下载图片失败 {image_url}: {e}")
            return None
    
    def _insert_image_into_placeholder(self, slide, img_stream: BytesIO, placeholder_shape_id: int):
        """将图片插入占位符"""
        logger.debug(f"尝试将图片插入占位符 {placeholder_shape_id}")
        
        placeholder_shape = None
        for shape in slide.shapes:
            if shape.shape_id == placeholder_shape_id:
                placeholder_shape = shape
                print(f"✓ 找到图片占位符 {placeholder_shape_id} (名称: {shape.name})")
                break
        
        if not placeholder_shape:
            logger.error(f"✗ 未找到图片占位符 {placeholder_shape_id}")
            available_ids = [s.shape_id for s in slide.shapes]
            logger.error(f"  可用的形状ID: {available_ids}")
            return None
        
        try:
            img = Image.open(img_stream)
            img_width, img_height = img.size
            logger.debug(f"图片尺寸: {img_width} x {img_height}")
            
            placeholder_width = placeholder_shape.width
            placeholder_height = placeholder_shape.height
            
            # 计算缩放比例
            width_ratio = placeholder_width / img_width
            height_ratio = placeholder_height / img_height
            scaling_factor = min(width_ratio, height_ratio)
            
            pic_width = int(img_width * scaling_factor)
            pic_height = int(img_height * scaling_factor)
            
            # 居中位置
            left = placeholder_shape.left + (placeholder_width - pic_width) / 2
            top = placeholder_shape.top + (placeholder_height - pic_height) / 2
            
            img_stream.seek(0)
            added_picture = slide.shapes.add_picture(img_stream, left, top, pic_width, pic_height)
            print(f"✓ 图片已插入，缩放后尺寸: {pic_width/914400:.2f} x {pic_height/914400:.2f} 英寸")
            
            # 移除占位符
            sp = placeholder_shape._element
            sp.getparent().remove(sp)
            
            return added_picture
        except Exception as e:
            logger.error(f"✗ 插入图片失败: {e}")
            return None
    
    def create_slide(self, image_data: Dict, section_title: str = "图表"):
        """创建图片幻灯片"""
        self.slide_counter += 1
        image_url = image_data.get("url", "")
        description = image_data.get("alt", "")
        image_title = section_title
        
        print(f"\n{'#'*80}")
        print(f"# 创建第 {self.slide_counter} 页: 图片页")
        print(f"# 标题: {image_title}")
        print(f"# 图片URL: {image_url}")
        print(f"# 是否有描述: {'是' if description else '否'}")
        print(f"{'#'*80}")
        
        img_stream = self._download_image(image_url)
        if not img_stream:
            logger.warning(f"跳过图片幻灯片: {image_url}")
            return
        
        try:
            img = Image.open(img_stream)
            img_width, img_height = img.size
            print(f"图片尺寸: {img_width} x {img_height}")
        except Exception as e:
            logger.error(f"打开图片失败 {image_url}: {e}")
            return
        
        description = self.text_processor.remove_html_tags(description)
        
        # 选择布局
        if description:
            if img_width > img_height:
                layout_key = "IMAGE_TITLE_AND_DESCRIPTION_WIDE"
                print("使用横向图片布局（带描述）")
            else:
                layout_key = "IMAGE_TITLE_AND_DESCRIPTION_TALL"
                print("使用纵向图片布局（带描述）")
        else:
            layout_key = "IMAGE_ONLY"
            print("使用纯图片布局")
        
        slide_layout = self._get_slide_layout(layout_key)
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 记录所有形状信息
        self._log_slide_shapes(slide, "图片页")
        
        self._insert_image_into_placeholder(slide, img_stream, self.config.SHAPE_IDS["IMAGE_PLACEHOLDER"])
        self._add_text_with_auto_fit(
            slide, 
            self.config.SHAPE_IDS["IMAGE_TITLE"], 
            image_title, 
            font_type="title",
            is_title_page=False
        )
        
        if description:
            self._add_text_with_auto_fit(
                slide,
                self.config.SHAPE_IDS["IMAGE_DESCRIPTION"],
                description,
                font_type="content"
            )
        
        self._fill_empty_placeholders(slide)
        print(f"✓ 图片页创建完成")

class SubSectionSlideStrategy(SlideStrategy):
    """子章节详情页生成策略"""
    
    def create_slide(self, section_name: str, sub_section_content: List[Dict]):
        """创建子章节幻灯片"""
        num_items = len(sub_section_content)
        if num_items == 0:
            logger.debug(f"跳过空子章节: '{section_name}'")
            return
        
        self.slide_counter += 1
        print(f"\n{'#'*80}")
        print(f"# 创建第 {self.slide_counter} 页: 子章节页")
        print(f"# 标题: {section_name}")
        print(f"# 子项数: {num_items}")
        print(f"{'#'*80}")
        
        # 根据项目数选择布局
        layout_map = {
            1: "TEXT_ONLY_SMALL_TITLE",
            2: "SUBCHAPTER_2_ITEMS",
            3: "SUBCHAPTER_3_ITEMS",
            4: "SUBCHAPTER_4_ITEMS",
            5: "SUBCHAPTER_5_ITEMS",
        }
        
        layout_key = layout_map.get(num_items)
        if not layout_key:
            logger.warning(f"不支持 {num_items} 个项目的子章节")
            return
        
        print(f"使用布局: {layout_key}")
        
        slide_layout = self._get_slide_layout(layout_key)
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 记录所有形状信息
        self._log_slide_shapes(slide, "子章节页")
        
        self._add_text_with_auto_fit(
            slide,
            self.config.SHAPE_IDS["SUBCHAPTER_TITLE"],
            section_name,
            font_type="title",
            is_title_page=False
        )
        
        # 添加内容项
        if num_items == 1:
            detail = sub_section_content[0].get('detail', '')
            print(f"添加单项内容: {detail[:50]}...")
            self._add_text_with_auto_fit(
                slide,
                self.config.SHAPE_IDS["CONTENT_TEXT"],
                detail,
                font_type="content"
            )
        elif num_items == 5:
            # 5项特殊处理
            print("处理5项内容（特殊布局）")
            for i, item in enumerate(sub_section_content[:5]):
                num_id, text_id = self.config.SHAPE_IDS["SUBCHAPTER_ITEMS"][5][i]
                print(f"  项目 {i+1}: 编号形状={num_id}, 文本形状={text_id}")
                self._add_text_with_auto_fit(slide, num_id, str(i+1), font_type="small")
                combined_text = f"{item.get('summary', '')}\n{item.get('detail', '')}"
                self._add_text_with_auto_fit(slide, text_id, combined_text, font_type="content")
        else:
            # 2-4项通用处理
            shape_ids = self.config.SHAPE_IDS["SUBCHAPTER_ITEMS"][num_items]
            print(f"处理 {num_items} 项内容，形状ID: {shape_ids}")
            for i, (item, shape_id) in enumerate(zip(sub_section_content, shape_ids)):
                summary = item.get("summary", "")
                detail = item.get("detail", "")
                combined_text = f"{summary}\n{detail}" if summary else detail
                print(f"  项目 {i+1}: 形状={shape_id}, 内容长度={len(combined_text)}")
                self._add_text_with_auto_fit(slide, shape_id, combined_text, font_type="content")
        
        self._fill_empty_placeholders(slide)
        print(f"✓ 子章节页创建完成")

class ReferencesSlideStrategy(SlideStrategy):
    """参考文献页生成策略"""
    
    def _process_reference_text(self, reference_text: str) -> str:
        """处理参考文献文本"""
        if not isinstance(reference_text, str):
            return ""
        lines = reference_text.splitlines()
        processed_lines = []
        i = 0
        while i < len(lines):
            current_line = lines[i].strip()
            if len(current_line) < 15 and i + 1 < len(lines):
                current_line += " " + lines[i + 1].strip()
                i += 1
            processed_lines.append(current_line)
            i += 1
        return "\n".join(processed_lines)
    
    def create_slide(self, references: List[str]):
        """创建参考文献幻灯片"""
        if not references:
            print("无参考文献，跳过")
            return
        
        print(f"总参考文献数: {len(references)}")
        
        # 预处理参考文献
        processed_references = [
            self._process_reference_text(self.text_processor.remove_html_tags(ref))
            for ref in references[:self.config.MAX_TOTAL_REFERENCES]
        ]
        
        # 分页处理
        total_pages = (len(processed_references) + self.config.MAX_REFERENCES_PER_SLIDE - 1) // self.config.MAX_REFERENCES_PER_SLIDE
        
        for page_idx in range(0, len(processed_references), self.config.MAX_REFERENCES_PER_SLIDE):
            self.slide_counter += 1
            current_page = page_idx // self.config.MAX_REFERENCES_PER_SLIDE + 1
            
            print(f"\n{'#'*80}")
            print(f"# 创建第 {self.slide_counter} 页: 参考文献页 ({current_page}/{total_pages})")
            print(f"{'#'*80}")
            
            group = processed_references[page_idx:page_idx + self.config.MAX_REFERENCES_PER_SLIDE]
            slide_layout = self._get_slide_layout("REFERENCES_PAGE")
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # 记录所有形状信息
            self._log_slide_shapes(slide, "参考文献页")
            
            self._add_text_with_auto_fit(
                slide,
                self.config.SHAPE_IDS["REFERENCES_TITLE"],
                "部分参考文献",
                font_type="title",
                is_title_page=False
            )
            
            # 添加参考文献
            for idx, ref_text in enumerate(group):
                if idx < len(self.config.SHAPE_IDS["REFERENCES"]):
                    ref_config = self.config.SHAPE_IDS["REFERENCES"][idx]
                    ref_num = page_idx + idx + 1
                    print(f"添加参考文献 {ref_num}: 编号形状={ref_config['num']}, 文本形状={ref_config['text']}")
                    
                    self._add_text_with_auto_fit(
                        slide,
                        ref_config["num"],
                        f"{ref_num}.",
                        font_type="small"
                    )
                    self._add_text_with_auto_fit(
                        slide,
                        ref_config["text"],
                        ref_text,
                        font_type="small"
                    )
            
            self._fill_empty_placeholders(slide)
            print(f"✓ 参考文献页创建完成")

class EndSlideStrategy(SlideStrategy):
    """结束页生成策略"""
    
    def create_slide(self):
        """创建结束页"""
        self.slide_counter += 1
        print(f"\n{'#'*80}")
        print(f"# 创建第 {self.slide_counter} 页: 结束页")
        print(f"{'#'*80}")
        
        slide_layout = self._get_slide_layout("END_PAGE")
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 记录所有形状信息
        self._log_slide_shapes(slide, "结束页")
        
        print(f"✓ 结束页创建完成")

# ==================== 主生成器类修改 ====================
class PresentationGenerator:
    """演示文稿生成器主类"""
    
    def __init__(self, template_file_name: str = 'ppt_template_0717.pptx'):
        self.current_dir = os.path.dirname(os.path.abspath(__file__))
        template_path = os.path.join(self.current_dir, template_file_name)
        
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"模板文件未找到: {template_path}")
        
        logger.info(f"\n{'='*80}")
        logger.info(f"初始化 PresentationGenerator")
        logger.info(f"模板文件: {template_path}")
        logger.info(f"{'='*80}")
        
        self.presentation = Presentation(template_path)
        self.config = SlideConfig()
        self.text_processor = TextProcessor()
        
        # 记录模板信息
        logger.info(f"模板布局数量: {len(self.presentation.slide_layouts)}")
        for idx, layout in enumerate(self.presentation.slide_layouts):
            logger.info(f"  布局 {idx}: {layout.name}")
        
        # 共享页面计数器
        self.slide_counter = 0
        
        # 初始化各种策略
        self.strategies = {
            "title": TitleSlideStrategy(self.presentation, self.config),
            "content": ContentSlideStrategy(self.presentation, self.config),
            "toc": TableOfContentsSlideStrategy(self.presentation, self.config),
            "image": ImageSlideStrategy(self.presentation, self.config),
            "subsection": SubSectionSlideStrategy(self.presentation, self.config),
            "references": ReferencesSlideStrategy(self.presentation, self.config),
            "end": EndSlideStrategy(self.presentation, self.config),
        }
        
        # 同步页面计数器
        for strategy in self.strategies.values():
            strategy.slide_counter = 0
        
        logger.info("初始化完成")
    
    def _parse_content_blocks(self, content_blocks: List[Dict]) -> Tuple[str, str, List[Dict]]:
        """解析内容块"""
        slide_title = ""
        main_paragraph_texts = []
        bullet_points_list = []
        
        logger.debug("开始解析内容块")
        
        for block in content_blocks:
            block_type = block.get("type")
            children = block.get("children", [])
            
            if block_type == "h1" and not slide_title:
                text_content = "".join(c.get("text", "") for c in children if c.get("text"))
                slide_title = self.text_processor.remove_html_tags(text_content)
                logger.debug(f"找到标题: {slide_title}")
            elif block_type == "p":
                text_content = "".join(c.get("text", "") for c in children if c.get("text"))
                clean_text = self.text_processor.remove_html_tags(text_content)
                if clean_text:
                    main_paragraph_texts.append(clean_text)
                    logger.debug(f"找到段落: {clean_text[:50]}...")
            elif block_type == "bullets":
                logger.debug(f"找到项目列表，包含 {len(children)} 项")
                for bullet_item in children:
                    if bullet_item.get("type") == "bullet":
                        summary_text = ""
                        detail_text = ""
                        for bullet_child in bullet_item.get("children", []):
                            if bullet_child.get("type") == "h3":
                                summary_text = "".join(
                                    c.get("text", "") for c in bullet_child.get("children", []) 
                                    if c.get("text"))
                            elif bullet_child.get("type") == "p":
                                detail_text = "".join(
                                    c.get("text", "") for c in bullet_child.get("children", [])
                                    if c.get("text"))
                        
                        bullet_points_list.append({
                            "summary": self.text_processor.remove_html_tags(summary_text),
                            "detail": self.text_processor.remove_html_tags(detail_text)
                        })
        
        logger.debug(f"解析完成: 标题='{slide_title}', 段落数={len(main_paragraph_texts)}, 项目数={len(bullet_points_list)}")
        
        return slide_title, "\n".join(main_paragraph_texts), bullet_points_list
    
    def _format_bullet_points_as_text(self, bullet_points: List[Dict]) -> str:
        """将bullet points格式化为文本内容"""
        formatted_text = []
        for i, item in enumerate(bullet_points):
            summary = item.get("summary", "")
            detail = item.get("detail", "")
            
            if summary:
                formatted_text.append(f"{i+1}. {summary}")
            if detail:
                formatted_text.append(f"   {detail}")
            formatted_text.append("")  # 添加空行分隔
        
        return "\n".join(formatted_text).strip()
    
    def generate_presentation(self, json_data: Dict) -> Optional[str]:
        """生成演示文稿的主方法"""
        logger.info(f"\n{'*'*80}")
        logger.info("开始生成演示文稿")
        logger.info(f"{'*'*80}")
        
        if not isinstance(json_data, dict):
            logger.error("无效输入：需要字典类型")
            return None
        
        # 获取文档标题
        doc_title = json_data.get("title", "")
        overall_references = json_data.get("references", [])
        sections = json_data.get("sections", [])
        
        logger.info(f"文档标题: '{doc_title}'")
        logger.info(f"章节数: {len(sections)}")
        logger.info(f"参考文献数: {len(overall_references)}")
        
        # 如果没有指定标题，从第一个section获取
        if not doc_title and sections:
            first_section_content = sections[0].get("content", [])
            for block in first_section_content:
                if block.get("type") == "h1":
                    doc_title = "".join(c.get("text", "") for c in block.get("children", []) if c.get("text"))
                    break
        
        if not doc_title:
            doc_title = "未命名演示文稿"
        
        try:
            # 重置所有策略的计数器
            for strategy in self.strategies.values():
                strategy.slide_counter = 0
            
            # 1. 创建标题页
            logger.info(f"\n{'='*60}")
            logger.info("创建标题页")
            logger.info(f"{'='*60}")
            self.strategies["title"].create_slide(doc_title)
            current_slide_count = 1
            for strategy in self.strategies.values():
                strategy.slide_counter = current_slide_count
            
            # 2. 处理每个section
            for section_idx, section_obj in enumerate(sections):
                logger.info(f"\n{'='*60}")
                logger.info(f"处理第 {section_idx + 1} 个章节")
                logger.info(f"{'='*60}")

                if section_idx == 0:
                    logger.info("跳过第一个章节的内容生成")
                    continue
                
                section_content = section_obj.get("content", [])
                section_root_image = section_obj.get("rootImage", {})
                
                # 解析内容
                slide_title, main_text, bullet_points = self._parse_content_blocks(section_content)
                
                # 如果这是第一个section且标题与文档标题相同
                if section_idx == 0 and slide_title == doc_title:
                    slide_title = "概述"
                
                # 决定使用哪种策略
                # 第一个section（第2页）：使用内容页
                if section_idx == 0:
                    if bullet_points:
                        bullet_text = self._format_bullet_points_as_text(bullet_points)
                        combined_text = main_text + "\n\n" + bullet_text if main_text else bullet_text
                        logger.info(f"创建内容页（包含{len(bullet_points)}个要点）")
                        self.strategies["content"].create_slide(slide_title, combined_text)
                        chunks = self.text_processor.split_text_into_chunks(combined_text)
                        current_slide_count += len(chunks)
                    elif main_text:
                        logger.info("创建内容页（纯文本）")
                        self.strategies["content"].create_slide(slide_title, main_text)
                        chunks = self.text_processor.split_text_into_chunks(main_text)
                        current_slide_count += len(chunks)
                
                # 第二个section（第3页）：使用SUBCHAPTER_3_ITEMS布局
                elif section_idx == 1:
                    if bullet_points and len(bullet_points) == 3:
                        logger.info("创建子章节页（使用SUBCHAPTER_3_ITEMS布局，ID=16）")
                        # 如果还有额外的段落文本，将其添加到最后一个bullet point
                        if main_text:
                            bullet_points[-1]['detail'] += f"\n\n{main_text}"
                        self.strategies["subsection"].create_slide(slide_title, bullet_points)
                        current_slide_count += 1
                    else:
                        # 如果不是正好3个bullet points，仍使用内容页
                        logger.info(f"注意：第二个section有{len(bullet_points)}个要点，不是3个，将使用普通内容页")
                        if bullet_points:
                            bullet_text = self._format_bullet_points_as_text(bullet_points)
                            combined_text = main_text + "\n\n" + bullet_text if main_text else bullet_text
                        else:
                            combined_text = main_text
                        self.strategies["content"].create_slide(slide_title, combined_text)
                        chunks = self.text_processor.split_text_into_chunks(combined_text)
                        current_slide_count += len(chunks)
                
                # 第三个section（第5页）：使用SUBCHAPTER_3_ITEMS布局
                elif section_idx == 2:
                    # 将段落转换为3个子项目格式
                    paragraphs = []
                    for block in section_content:
                        if block.get("type") == "p":
                            text_content = "".join(c.get("text", "") for c in block.get("children", []) if c.get("text"))
                            if text_content:
                                paragraphs.append(text_content)
                    
                    if len(paragraphs) >= 2:
                        # 创建3个子项目，如果只有2个段落，第三个留空或合并
                        sub_items = []
                        if len(paragraphs) == 2:
                            # 将第二个段落分成两部分
                            second_para_sentences = paragraphs[1].split('. ')
                            if len(second_para_sentences) >= 2:
                                mid_point = len(second_para_sentences) // 2
                                para2_part1 = '. '.join(second_para_sentences[:mid_point]) + '.'
                                para2_part2 = '. '.join(second_para_sentences[mid_point:])
                                
                                sub_items = [
                                    {"summary": "Research Direction", "detail": paragraphs[0]},
                                    {"summary": "Combination Strategies", "detail": para2_part1},
                                    {"summary": "Biomarker Development", "detail": para2_part2}
                                ]
                            else:
                                sub_items = [
                                    {"summary": "Research Direction", "detail": paragraphs[0]},
                                    {"summary": "Future Studies", "detail": paragraphs[1]},
                                    {"summary": "", "detail": ""}  # 空项
                                ]
                        
                        logger.info("创建子章节页（使用SUBCHAPTER_3_ITEMS布局，ID=16）")
                        self.strategies["subsection"].create_slide(slide_title, sub_items)
                        current_slide_count += 1
                    else:
                        # 如果内容不适合，使用普通内容页
                        logger.info("内容不适合SUBCHAPTER_3_ITEMS布局，使用普通内容页")
                        self.strategies["content"].create_slide(slide_title, main_text)
                        chunks = self.text_processor.split_text_into_chunks(main_text)
                        current_slide_count += len(chunks)
                
                # 同步计数器
                for strategy in self.strategies.values():
                    strategy.slide_counter = current_slide_count
                
                # 处理图片（非背景图）- 第4页
                if section_root_image and section_root_image.get("url") and not section_root_image.get("background", False):
                    logger.info("创建图片页")
                    self.strategies["image"].create_slide(section_root_image, slide_title)
                    current_slide_count += 1
                    for strategy in self.strategies.values():
                        strategy.slide_counter = current_slide_count
            
            # 3. 添加参考文献页
            if overall_references:
                logger.info(f"\n{'='*60}")
                logger.info("创建参考文献页")
                logger.info(f"{'='*60}")
                pages_needed = (len(overall_references[:SlideConfig.MAX_TOTAL_REFERENCES]) + 
                            SlideConfig.MAX_REFERENCES_PER_SLIDE - 1) // SlideConfig.MAX_REFERENCES_PER_SLIDE
                self.strategies["references"].create_slide(overall_references)
                current_slide_count += pages_needed
                for strategy in self.strategies.values():
                    strategy.slide_counter = current_slide_count
            
            # 4. 添加结束页
            logger.info(f"\n{'='*60}")
            logger.info("创建结束页")
            logger.info(f"{'='*60}")
            self.strategies["end"].create_slide()
            
            # 保存文件
            output_dir = os.path.join(self.current_dir, 'output_ppts')
            os.makedirs(output_dir, exist_ok=True)
            
            sanitized_title = re.sub(r'[\\/:*?"<>|]', '', doc_title)
            output_filename = os.path.join(output_dir, f'{sanitized_title}.pptx')
            self.presentation.save(output_filename)
            
            logger.info(f"\n{'*'*80}")
            logger.info(f"PPT生成成功!")
            logger.info(f"文件路径: {output_filename}")
            logger.info(f"总页数: {self.strategies['end'].slide_counter}")
            logger.info(f"{'*'*80}")
            
            # 输出页面总结
            logger.info("\n页面结构总结：")
            logger.info("第1页: 标题页")
            logger.info("第2页: 内容页 - 概述（3个要点）")
            logger.info("第3页: 子章节页（3项，使用布局ID=16）")
            logger.info("第4页: 图片页 - ")
            logger.info("第5页: 子章节页（3项，使用布局ID=16）")
            logger.info("第6页: 参考文献页")
            logger.info("第7页: 结束页")
            
            return output_filename
            
        except Exception as e:
            logger.critical(f"PPT生成失败: {e}", exc_info=True)
            return None

# ==================== 入口函数 ====================
def start_generate_presentation(json_input: Any) -> Optional[str]:
    """PPT生成的入口函数"""
    print("\n" + "="*100)
    print("PPT生成器启动")
    print("="*100)
    
    try:
        json_data = json.loads(json_input) if isinstance(json_input, str) else json_input
    except (json.JSONDecodeError, TypeError) as e:
        logger.error(f"JSON解析失败: {e}")
        return None
    
    generator = PresentationGenerator()
    output_path = generator.generate_presentation(json_data)
    
    return output_path

# 使用示例
if __name__ == "__main__":
    # 测试数据
    ppt_data_wrapper = [
    {
        "id": "8nJVwBp86wBXHnjf4n6Pl",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "Introduction to Tesla Motors"
                    }
                ],
                "id": "ZigbEk0pDa"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "研究主题"
                                    }
                                ],
                                "id": "SvHtGLhhhA"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉汽车公司的背景、历史和其在汽车行业的意义。"
                                    }
                                ],
                                "id": "giy9t5HNCg"
                            }
                        ],
                        "id": "TEkXhPZR3H"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "子问题拆解"
                                    }
                                ],
                                "id": "rTL7eI0msg"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "包括特斯拉的起源与发展历程、技术突破与市场地位、对传统汽车工业的影响以及未来的发展潜力与挑战。"
                                    }
                                ],
                                "id": "uhRwtMK5sW"
                            }
                        ],
                        "id": "4bOUOgNEQW"
                    }
                ],
                "id": "42w6LOyNtd"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Tesla car",
            "url": "https://cdn.pixabay.com/photo/2024/12/18/07/57/aura-9274671_640.jpg",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    },
    {
        "id": "ZBZE7AG3vQEucmtTqsgAW",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "特斯拉的市场地位与竞争分析"
                    }
                ],
                "id": "iJpwM0hq3o"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "市场份额"
                                    }
                                ],
                                "id": "qup7hHjvjs"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "根据2023年的数据，特斯拉在全球电动车市场的占有率约为15%，位居前列。"
                                    }
                                ],
                                "id": "JHh50y9f0R"
                            }
                        ],
                        "id": "iXhGmjEWjz"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "主要竞争对手"
                                    }
                                ],
                                "id": "dqCoOJbxSX"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "比亚迪、蔚来、宝马和大众等品牌在电动车市场中占据重要位置。"
                                    }
                                ],
                                "id": "SJ3fZdgs4_"
                            }
                        ],
                        "id": "VO1ifHwTbK"
                    }
                ],
                "id": "DlJf720Q5X"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Tesla market position",
            "url": "https://c-ssl.duitang.com/uploads/blog/202111/27/20211127170036_ecc10.png",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    },
    {
        "id": "Uvao5htlAoNDrrPAb0-DP",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "特斯拉的技术突破与市场地位"
                    }
                ],
                "id": "cDAtjVv-GY"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "技术突破"
                                    }
                                ],
                                "id": "Kdxr0Y7szM"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉在电池技术、自动驾驶系统（Autopilot）和软件更新方面具有显著优势。"
                                    }
                                ],
                                "id": "6IwSQCvPOw"
                            }
                        ],
                        "id": "EUcuywljYH"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "市场地位"
                                    }
                                ],
                                "id": "kL2OKt_WBY"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉在全球电动车市场中占据领先地位，凭借技术创新、品牌影响力和垂直整合模式建立了显著的竞争优势。"
                                    }
                                ],
                                "id": "090kJjLVcB"
                            }
                        ],
                        "id": "Uu6AxAwSwu"
                    }
                ],
                "id": "AE2uQeh6hp"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Tesla technology and market position",
            "url": "https://n.sinaimg.cn/spider20250621/120/w1440h1080/20250621/f8d9-7d234d7b43fda7ec916d01fb81555bae.jpg",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    },
    {
        "id": "RH_pGO5PnCWXzNXNoQmDI",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "特斯拉的可持续发展与环境影响"
                    }
                ],
                "id": "l1riZ-6GSg"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "减少碳足迹"
                                    }
                                ],
                                "id": "4fhg4XV0l5"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉通过制造电动汽车减少了传统燃油车的碳排放，尤其是在电力来源清洁的地区，其全生命周期碳排放显著低于燃油车。"
                                    }
                                ],
                                "id": "DOK4Ws-nI4"
                            }
                        ],
                        "id": "EqWLFcVAjn"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "推动可再生能源"
                                    }
                                ],
                                "id": "ugDNFCJwin"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉开发了太阳能产品和储能系统，鼓励用户采用清洁能源。公司致力于建设全球范围内的超级充电站，并逐步转向使用可再生能源供电。"
                                    }
                                ],
                                "id": "jAPWqlhq00"
                            }
                        ],
                        "id": "CxW3-jDihT"
                    }
                ],
                "id": "CmPoEt9TBE"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Battery degradation curve",
            "url": "https://n.sinaimg.cn/spider20250610/598/w955h443/20250610/d5d4-b0eb5848953ded18fbc44aefb452923d.png",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    },
    {
        "id": "5uyW3JgOeD9ZbEFB_2iW_",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "特斯拉面临的挑战与未来前景"
                    }
                ],
                "id": "y2QFM73QyV"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "主要挑战"
                                    }
                                ],
                                "id": "q1dtux9SFi"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉面临的主要生产瓶颈包括工厂产能受限、自动化生产线故障频发，以及供应链问题如芯片短缺和电池材料供应不足。"
                                    }
                                ],
                                "id": "M2egHQqaSe"
                            }
                        ],
                        "id": "cUdLtk1Ahm"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "未来增长潜力"
                                    }
                                ],
                                "id": "0CVphgBEZr"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉在全球市场扩张中具有巨大潜力，特别是在新兴市场如印度和东南亚。此外，FSD（完全自动驾驶）技术和电池技术的突破将为其带来新的收入来源。"
                                    }
                                ],
                                "id": "MXsd1dQAI-"
                            }
                        ],
                        "id": "6ffe1enX75"
                    }
                ],
                "id": "x1EUFy0C3R"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Tesla challenges and future prospects",
            "url": "https://n.sinaimg.cn/spider20250610/598/w955h443/20250610/d5d4-b0eb5848953ded18fbc44aefb452923d.png",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    },
    {
        "id": "kp8wqwwK5wj_aZIg8OsC4",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "特斯拉的市场扩张与全球化战略"
                    }
                ],
                "id": "g-QwuNImsJ"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "全球市场扩张"
                                    }
                                ],
                                "id": "UEpuyIdugA"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉正在推进本地化生产计划，例如在中国和德国建设超级工厂，以降低运输成本并提高市场响应速度。"
                                    }
                                ],
                                "id": "Xb3pXGDIaz"
                            }
                        ],
                        "id": "lBazhK7w1m"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "新兴市场机会"
                                    }
                                ],
                                "id": "galszc55Mg"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "印度、东南亚等新兴市场成为特斯拉的重点扩展方向，但需克服政策和基础设施限制。"
                                    }
                                ],
                                "id": "wXw0AdR7yl"
                            }
                        ],
                        "id": "yyBao4j3YB"
                    }
                ],
                "id": "b9J1xaY5X9"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Tesla global expansion strategy",
            "url": "https://n.sinaimg.cn/spider20250610/598/w955h443/20250610/d5d4-b0eb5848953ded18fbc44aefb452923d.png",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    },
    {
        "id": "7voQ5e4hggtoaSek_YGGw",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "特斯拉的全球化战略与本地化生产"
                    }
                ],
                "id": "Lux82bMEPV"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "全球市场扩张"
                                    }
                                ],
                                "id": "P0fspVVGr0"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉正在推进本地化生产计划，例如在中国和德国建设超级工厂，以降低运输成本并提高市场响应速度。"
                                    }
                                ],
                                "id": "Jpxrqh_qGG"
                            }
                        ],
                        "id": "mwa91YKrbu"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "新兴市场机会"
                                    }
                                ],
                                "id": "80vK3tTymo"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "印度、东南亚等新兴市场成为特斯拉的重点扩展方向，但需克服政策和基础设施限制。"
                                    }
                                ],
                                "id": "8tSM2IFQOL"
                            }
                        ],
                        "id": "z7GH6KrCKs"
                    }
                ],
                "id": "c463mxyNyi"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Tesla global expansion strategy",
            "url": "https://n.sinaimg.cn/spider20250610/598/w955h443/20250610/d5d4-b0eb5848953ded18fbc44aefb452923d.png",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    },
    {
        "id": "dsTWMJpHAUmdLQ2bQAwVx",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "特斯拉的全球化战略与本地化生产"
                    }
                ],
                "id": "TmVVrZawhs"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "全球市场扩张"
                                    }
                                ],
                                "id": "B5Uphjvhyw"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉正在推进本地化生产计划，例如在中国和德国建设超级工厂，以降低运输成本并提高市场响应速度。"
                                    }
                                ],
                                "id": "g_sd3C2Txv"
                            }
                        ],
                        "id": "QsUaSuQiVI"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "新兴市场机会"
                                    }
                                ],
                                "id": "PQkyvQuyQk"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "印度、东南亚等新兴市场成为特斯拉的重点扩展方向，但需克服政策和基础设施限制。"
                                    }
                                ],
                                "id": "EbHbBGaBzS"
                            }
                        ],
                        "id": "UtFU5zml4z"
                    }
                ],
                "id": "n4NVulKJzn"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Tesla global expansion strategy",
            "url": "https://cdn.pixabay.com/photo/2024/12/18/15/02/old-9275581_640.jpg",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    },
    {
        "id": "zfWnbseSqZveQtxH5KHTG",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "特斯拉的全球市场扩张与本地化生产"
                    }
                ],
                "id": "8OTGC3IPFd"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "全球市场扩张"
                                    }
                                ],
                                "id": "fJ_qZs3b4-"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉正在推进本地化生产计划，例如在中国和德国建设超级工厂，以降低运输成本并提高市场响应速度。"
                                    }
                                ],
                                "id": "iTMHgRX0Ur"
                            }
                        ],
                        "id": "5zbw3jLmCt"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "新兴市场机会"
                                    }
                                ],
                                "id": "XMbrJ7khQy"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "印度、东南亚等新兴市场成为特斯拉的重点扩展方向，但需克服政策和基础设施限制。"
                                    }
                                ],
                                "id": "lWuTdDTpMr"
                            }
                        ],
                        "id": "--E-W3gMGT"
                    }
                ],
                "id": "E8rt3Dcr28"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Tesla global expansion strategy",
            "url": "https://n.sinaimg.cn/spider20250610/598/w955h443/20250610/d5d4-b0eb5848953ded18fbc44aefb452923d.png",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    },
    {
        "id": "TEqcpxtIbejZAEY47cWv0",
        "content": [
            {
                "type": "h1",
                "children": [
                    {
                        "text": "特斯拉的可持续发展与环境影响"
                    }
                ],
                "id": "4zF69jZYen"
            },
            {
                "type": "bullets",
                "children": [
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "减少碳足迹"
                                    }
                                ],
                                "id": "R8ChCJYkVK"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉通过制造电动汽车减少了传统燃油车的碳排放，尤其是在电力来源清洁的地区，其全生命周期碳排放显著低于燃油车。"
                                    }
                                ],
                                "id": "DRqPQzvhrc"
                            }
                        ],
                        "id": "ey6nZ58uu7"
                    },
                    {
                        "type": "bullet",
                        "children": [
                            {
                                "type": "h3",
                                "children": [
                                    {
                                        "text": "推动可再生能源"
                                    }
                                ],
                                "id": "AzBGurXru0"
                            },
                            {
                                "type": "p",
                                "children": [
                                    {
                                        "text": "特斯拉开发了太阳能产品和储能系统，鼓励用户采用清洁能源。公司致力于建设全球范围内的超级充电站，并逐步转向使用可再生能源供电。"
                                    }
                                ],
                                "id": "-wcBBI4GqI"
                            }
                        ],
                        "id": "eYUyIUhhsB"
                    }
                ],
                "id": "dlyHYheu0A"
            }
        ],
        "alignment": "center",
        "rootImage": {
            "alt": "Battery degradation curve",
            "url": "https://n.sinaimg.cn/spider20250610/598/w955h443/20250610/d5d4-b0eb5848953ded18fbc44aefb452923d.png",
            "query": "",
            "background": False
        },
        "layoutType": "vertical"
    }
]

    output_file_path = start_generate_presentation(ppt_data_wrapper)
    if output_file_path:
        print(f"\nPPT生成完成: {output_file_path}")
    else:
        print("\nPPT生成失败")