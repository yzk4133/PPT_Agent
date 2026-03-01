"""
LangChain Native Tool: Create PowerPoint

Implements PowerPoint generation functionality using python-pptx.
"""

import logging
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict

from pptx import Presentation
from pptx.util import Inches
from langchain_core.tools import StructuredTool
from langchain_core.pydantic_v1 import BaseModel, Field

from backend.tools.core.monitoring import monitor_tool
from backend.tools.application.tool_registry import get_native_registry

logger = logging.getLogger(__name__)


# Input schema
class CreatePptxInput(BaseModel):
    """PowerPoint creation input schema"""

    slides: List[Dict] = Field(description="List of slide data dictionaries")
    output_path: Optional[str] = Field(default=None, description="Output file path (optional)")
    template_path: Optional[str] = Field(default=None, description="Template file path (optional)")


@monitor_tool
async def create_pptx(
    slides: List[Dict], output_path: Optional[str] = None, template_path: Optional[str] = None
) -> dict:
    """
    Create PowerPoint file from structured data

    Generates a PowerPoint presentation from slide data including titles,
    content, images, and notes.

    Args:
        slides: List of slide dictionaries with title, content, images
        output_path: Output file path (auto-generated if None)
        template_path: Optional PowerPoint template file

    Returns:
        Dictionary with output path and metadata

    Example slide:
    {
        "layout": "Title and Content",
        "title": "Slide Title",
        "subtitle": "Optional subtitle",
        "content": ["Point 1", "Point 2"],
        "images": ["url1"],
        "notes": "Speaker notes"
    }
    """
    logger.info(f"[create_pptx] Creating PPT with {len(slides)} slides")

    try:
        # Generate output path if not provided
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(Path(tempfile.gettempdir()) / f"presentation_{timestamp}.pptx")

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Load template or create new presentation
        if template_path and Path(template_path).exists():
            prs = Presentation(template_path)
            logger.info(f"[create_pptx] Loaded template: {template_path}")
        else:
            prs = Presentation()

        # Add slides
        for slide_data in slides:
            _add_slide(prs, slide_data)

        # Save presentation
        prs.save(str(output_path))

        # Get file size
        file_size = output_path.stat().st_size / (1024 * 1024)  # MB

        logger.info(f"[create_pptx] Created {output_path} ({file_size:.2f} MB)")

        return {
            "output_path": str(output_path),
            "total_slides": len(slides),
            "file_size_mb": round(file_size, 2),
            "template_used": template_path is not None,
        }

    except Exception as e:
        logger.error(f"[create_pptx] Error: {e}", exc_info=True)
        raise


def _add_slide(prs: Presentation, slide_data: Dict):
    """Add a single slide to the presentation"""
    layout_name = slide_data.get("layout", "Title and Content")
    layout = _get_layout(prs, layout_name)

    if layout:
        slide = prs.slides.add_slide(layout)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Set title
    title = slide_data.get("title", "")
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Set subtitle if present
    if "subtitle" in slide_data:
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # Subtitle type
                shape.text = slide_data["subtitle"]
                break

    # Add content
    content = slide_data.get("content", [])
    if content and layout_name != "Blank":
        _add_content(slide, content)

    # Add images
    images = slide_data.get("images", [])
    for img_data in images:
        _add_image(slide, img_data)

    # Add notes
    if "notes" in slide_data:
        slide.notes_slide.notes_text_frame.text = slide_data["notes"]


def _add_content(slide, content: List[str]):
    """Add text content to slide"""
    # Try to find content placeholder
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 1:  # Body type
            content_placeholder = shape
            break

    if not content_placeholder:
        return

    text_frame = content_placeholder.text_frame
    text_frame.clear()  # Clear default text

    # Add first paragraph
    if content:
        text_frame.text = str(content[0])

        # Add remaining paragraphs
        for point in content[1:]:
            p = text_frame.add_paragraph()
            p.text = str(point)
            p.level = 0


def _add_image(slide, img_data):
    """Add image to slide"""
    # Get image URL or path
    if isinstance(img_data, str):
        img_url = img_url = img_data
        img_x, img_y, img_w, img_h = 1, 1, 3, 2
    elif isinstance(img_data, dict):
        img_url = img_data.get("url", "")
        img_x = img_data.get("x", 1)
        img_y = img_data.get("y", 1)
        img_w = img_data.get("width", 3)
        img_h = img_data.get("height", 2)
    else:
        return

    # Check if it's a URL or local file
    if img_url.startswith(("http://", "https://")):
        # For URL, add a placeholder text box
        left = Inches(img_x)
        top = Inches(img_y)
        width = Inches(img_w)
        height = Inches(img_h)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = f"[Image: {img_url}]"
        text_frame.word_wrap = True
    else:
        # Local file
        try:
            slide.shapes.add_picture(
                img_url,
                left=Inches(img_x),
                top=Inches(img_y),
                width=Inches(img_w),
                height=Inches(img_h),
            )
        except Exception as e:
            logger.warning(f"[create_pptx] Failed to add image {img_url}: {e}")


def _get_layout(prs: Presentation, layout_name: str):
    """Get layout by name"""
    layout_map = {
        "Title": 0,
        "Title and Content": 1,
        "Section Header": 2,
        "Two Content": 3,
        "Comparison": 4,
        "Title Only": 5,
        "Blank": 6,
    }
    idx = layout_map.get(layout_name, 1)
    if idx < len(prs.slide_layouts):
        return prs.slide_layouts[idx]
    return prs.slide_layouts[0]


# Create LangChain StructuredTool
tool = StructuredTool.from_function(
    func=create_pptx,
    name="create_pptx",
    description="Create PowerPoint files from structured data. Use this to generate PPT presentations.",
    args_schema=CreatePptxInput,
)

# Auto-register with global registry
registry = get_native_registry()
registry.register_tool(tool, category=registry.UTILITY)

__all__ = ["tool", "create_pptx"]
