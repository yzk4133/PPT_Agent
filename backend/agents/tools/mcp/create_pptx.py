#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
MCP Tool: Create PPT

Implements PowerPoint generation functionality using python-pptx.
"""

import os
import json
import logging
from typing import Optional, Dict, Any, List
from datetime import datetime
from pathlib import Path
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from .base_mcp_tool import BaseMCPTool

logger = logging.getLogger(__name__)

class CreatePptxTool(BaseMCPTool):
    """PowerPoint file creation tool using python-pptx"""

    def __init__(self):
        super().__init__(
            name="create_pptx",
            description="Create PowerPoint files from structured data"
        )

    async def execute(
        self,
        slides: List[Dict],
        output_path: Optional[str] = None,
        template_path: Optional[str] = None,
        theme: Optional[Dict] = None,
        tool_context: Optional[Any] = None
    ) -> str:
        """
        Create PowerPoint file

        Args:
            slides: List of slide data dictionaries
            output_path: Output file path (if None, uses temp directory)
            template_path: Optional template file path
            theme: Optional theme configuration
            tool_context: Optional tool context

        Returns:
            JSON string with result
        """
        logger.info(f"[create_pptx] Creating PPT with {len(slides)} slides")

        try:
            # Generate output path if not provided
            if output_path is None:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = str(Path(tempfile.gettempdir()) / f"presentation_{timestamp}.pptx")

            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # Load template or create new presentation
            if template_path and Path(template_path).exists():
                prs = Presentation(template_path)
                logger.info(f"Loaded template: {template_path}")
            else:
                prs = Presentation()
                # Apply theme if provided
                if theme:
                    self._apply_theme(prs, theme)

            # Add slides
            for slide_data in slides:
                self._add_slide(prs, slide_data)

            # Save presentation
            prs.save(str(output_path))

            # Get file size
            file_size = output_path.stat().st_size / (1024 * 1024)  # MB

            metadata = {
                "output_path": str(output_path),
                "total_slides": len(slides),
                "file_size_mb": round(file_size, 2),
                "template_used": template_path is not None
            }

            return self._success(metadata, metadata=metadata)

        except Exception as e:
            logger.error(f"Create PPT error: {e}", exc_info=True)
            return self._error(
                message=str(e),
                code="CREATE_ERROR",
                details={"output_path": output_path}
            )

    def _add_slide(self, prs: Presentation, slide_data: Dict):
        """Add a single slide to the presentation"""
        layout_name = slide_data.get("layout", "Title and Content")
        layout = self._get_layout(prs, layout_name)

        # Remove default blank slide if needed
        if len(prs.slides) == 0 and layout_name == "Blank":
            slide = prs.slides.add_slide(layout)
        elif layout:
            slide = prs.slides.add_slide(layout)
        else:
            # Fallback to title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])

        # Set title
        title = slide_data.get("title", "")
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Set subtitle if present
        if "subtitle" in slide_data:
            # Try to find subtitle placeholder
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # Subtitle type
                    shape.text = slide_data["subtitle"]
                    break

        # Add content
        content = slide_data.get("content", [])
        if content and layout_name != "Blank":
            self._add_content(slide, content, slide_data)

        # Add images
        images = slide_data.get("images", [])
        for img_data in images:
            self._add_image(slide, img_data, slide_data)

        # Add notes
        if "notes" in slide_data:
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    def _add_content(self, slide, content: List[str], slide_data: Dict):
        """Add text content to slide"""
        # Try to find content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # Body type
                content_placeholder = shape
                break

        if not content_placeholder:
            return

        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Clear default text

        # Add first paragraph
        if content:
            text_frame.text = str(content[0])

            # Add remaining paragraphs
            for point in content[1:]:
                p = text_frame.add_paragraph()
                p.text = str(point)
                p.level = 0

    def _add_image(self, slide, img_data: Any, slide_data: Dict):
        """Add image to slide"""
        # Get image URL or path
        if isinstance(img_data, str):
            img_url = img_data
            img_x, img_y, img_w, img_h = 1, 1, 3, 2
        elif isinstance(img_data, dict):
            img_url = img_data.get("url", "")
            img_x = img_data.get("x", 1)
            img_y = img_data.get("y", 1)
            img_w = img_data.get("width", 3)
            img_h = img_data.get("height", 2)
        else:
            return

        # Check if it's a URL or local file
        if img_url.startswith(("http://", "https://")):
            # For URL, we'd need to download it first
            # For now, add a placeholder text box
            left = Inches(img_x)
            top = Inches(img_y)
            width = Inches(img_w)
            height = Inches(img_h)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = f"[Image: {img_url}]"
            text_frame.word_wrap = True
        else:
            # Local file
            try:
                slide.shapes.add_picture(
                    img_url,
                    left=Inches(img_x),
                    top=Inches(img_y),
                    width=Inches(img_w),
                    height=Inches(img_h)
                )
            except Exception as e:
                logger.warning(f"Failed to add image {img_url}: {e}")

    def _get_layout(self, prs: Presentation, layout_name: str):
        """Get layout by name"""
        layout_map = {
            "Title": 0,
            "Title and Content": 1,
            "Section Header": 2,
            "Two Content": 3,
            "Comparison": 4,
            "Title Only": 5,
            "Blank": 6
        }
        idx = layout_map.get(layout_name, 1)
        if idx < len(prs.slide_layouts):
            return prs.slide_layouts[idx]
        return prs.slide_layouts[0]

    def _apply_theme(self, prs: Presentation, theme: Dict):
        """Apply theme to presentation"""
        # Theme configuration
        # Note: python-pptx has limited theme support
        # This is a placeholder for future enhancement
        logger.info(f"Applying theme: {theme.get('name', 'custom')}")

# Global instance
_tool_instance = None

def get_tool() -> CreatePptxTool:
    """Get or create the create PPT tool instance"""
    global _tool_instance
    if _tool_instance is None:
        _tool_instance = CreatePptxTool()
    return _tool_instance

async def create_pptx(
    slides: List[Dict],
    output_path: Optional[str] = None,
    template_path: Optional[str] = None,
    theme: Optional[Dict] = None,
    tool_context: Optional[Any] = None
) -> str:
    """
    Create PowerPoint file

    Args:
        slides: List of slide data dictionaries
        output_path: Output file path (if None, uses temp directory)
        template_path: Optional template file path
        theme: Optional theme configuration
        tool_context: Optional tool context

    Returns:
        JSON string with result

    Example slide data:
    {
        "layout": "Title and Content",
        "title": "Slide Title",
        "subtitle": "Optional subtitle",
        "content": ["Point 1", "Point 2", "Point 3"],
        "images": ["url1", "url2"],
        "notes": "Speaker notes"
    }
    """
    tool = get_tool()
    return await tool.execute(
        slides=slides,
        output_path=output_path,
        template_path=template_path,
        theme=theme,
        tool_context=tool_context
    )
